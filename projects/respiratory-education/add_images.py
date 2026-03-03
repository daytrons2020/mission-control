#!/usr/bin/env python3
"""
Add images to PowerPoint presentations for Respiratory Education.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/root/.openclaw/workspace/projects/respiratory-education/output"
IMAGE_DIR = "/root/.openclaw/workspace/projects/respiratory-education/images"

def add_image_to_slide(prs_file, slide_idx, image_path, left, top, width, height):
    """Add an image to a specific slide."""
    if not os.path.exists(image_path):
        print(f"Image not found: {image_path}")
        return False
    
    prs = Presentation(prs_file)
    if slide_idx >= len(prs.slides):
        print(f"Slide index {slide_idx} out of range")
        return False
    
    slide = prs.slides[slide_idx]
    
    try:
        slide.shapes.add_picture(image_path, left, top, width, height)
        prs.save(prs_file)
        print(f"Added image to slide {slide_idx + 1} of {os.path.basename(prs_file)}")
        return True
    except Exception as e:
        print(f"Error adding image: {e}")
        return False

# Add images to 04-Patient-Assessment.pptx
print("Adding images to 04-Patient-Assessment.pptx...")
ppt_file = os.path.join(OUTPUT_DIR, "04-Patient-Assessment.pptx")

# Slide 12 (Physical Examination - Inspection) - add lung diagram
add_image_to_slide(ppt_file, 11, 
    os.path.join(IMAGE_DIR, "lung_diagram.png"),
    Inches(6.5), Inches(2), Inches(3), Inches(3))

# Slide 15 (Auscultation) - add stethoscope image
add_image_to_slide(ppt_file, 14,
    os.path.join(IMAGE_DIR, "stethoscope.jpg"),
    Inches(6.5), Inches(2), Inches(3), Inches(3))

# Slide 17 (Chest X-Ray) - add chest xray
add_image_to_slide(ppt_file, 16,
    os.path.join(IMAGE_DIR, "chest_xray.jpg"),
    Inches(6.5), Inches(2), Inches(3), Inches(3))

# Add images to 05-Emergency-Procedures.pptx
print("\nAdding images to 05-Emergency-Procedures.pptx...")
ppt_file = os.path.join(OUTPUT_DIR, "05-Emergency-Procedures.pptx")

# Slide 5 (BVM) - add BVM image
add_image_to_slide(ppt_file, 4,
    os.path.join(IMAGE_DIR, "bag_valve_mask.jpg"),
    Inches(6.5), Inches(2), Inches(3), Inches(3))

# Slide 8 (Intubation) - add intubation image
add_image_to_slide(ppt_file, 7,
    os.path.join(IMAGE_DIR, "intubation.jpg"),
    Inches(6.5), Inches(2), Inches(3), Inches(3))

# Add images to Mechanical-Ventilation-Basics.pptx
print("\nAdding images to Mechanical-Ventilation-Basics.pptx...")
ppt_file = os.path.join(OUTPUT_DIR, "Mechanical-Ventilation-Basics.pptx")

# Slide 8 (Ventilator Graphics) - add ventilator image
add_image_to_slide(ppt_file, 7,
    os.path.join(IMAGE_DIR, "ventilator.jpg"),
    Inches(6.5), Inches(2), Inches(3), Inches(3))

print("\nDone adding images!")
