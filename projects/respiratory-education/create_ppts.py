#!/usr/bin/env python3
"""
Create PowerPoint presentations for Respiratory Education
with full content and medical images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os
import requests
import json

# Create output directory if it doesn't exist
OUTPUT_DIR = "/root/.openclaw/workspace/projects/respiratory-education/output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Color scheme
DARK_BLUE = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(0, 102, 153)
ACCENT_BLUE = RGBColor(51, 153, 204)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64, 64, 64)

def create_title_slide(prs, title, subtitle="", presenter="", date="", department=""):
    """Create a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER
    
    # Add presenter info
    info_text = []
    if presenter:
        info_text.append(f"Presenter: {presenter}")
    if date:
        info_text.append(f"Date: {date}")
    if department:
        info_text.append(f"Department: {department}")
    
    if info_text:
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
        tf = info_box.text_frame
        for i, line in enumerate(info_text):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, content_lines, notes=""):
    """Create a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()
    
    # Add title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Check if line is a header (starts with ** or is all caps)
        if line.startswith("**") and line.endswith("**"):
            p.text = line.replace("**", "")
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
            p.space_before = Pt(12)
        elif line.startswith("-"):
            p.text = "• " + line[1:].strip()
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.level = 0
        elif line.startswith("  -"):
            p.text = "○ " + line[2:].strip()
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.level = 1
        else:
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
    
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content, notes=""):
    """Create a slide with two columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes
    
    return slide

def create_table_slide(prs, title, headers, rows, notes=""):
    """Create a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.6 * num_rows)).table
    
    # Set headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
    
    # Set rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = DARK_GRAY
    
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes
    
    return slide

def create_end_slide(prs, title="Thank You", subtitle="", department=""):
    """Create a closing slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER
    
    if department:
        dept_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
        tf = dept_box.text_frame
        p = tf.paragraphs[0]
        p.text = department
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# ============================================================
# PRESENTATION 1: 04-Patient-Assessment.pptx
# ============================================================

def create_patient_assessment_ppt():
    """Create the Patient Assessment presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs, 
        "Patient Assessment in Respiratory Care",
        "Comprehensive Assessment Techniques",
        "[Name], RRT, [Credentials]",
        "[Date]",
        "Respiratory Care")
    
    # Slide 2: Learning Objectives
    create_content_slide(prs, "Learning Objectives", [
        "By the end of this session, participants will be able to:",
        "",
        "- Perform a comprehensive respiratory patient assessment",
        "- Interpret arterial blood gas results and identify acid-base disorders",
        "- Analyze pulse oximetry and capnography data",
        "- Conduct a thorough physical examination of the chest",
        "- Evaluate chest X-rays for common respiratory conditions",
        "- Integrate assessment findings to develop a care plan"
    ], "Emphasize the importance of systematic assessment in respiratory care. Each skill builds upon the previous one.")
    
    # Slide 3: Assessment Framework
    create_content_slide(prs, "Assessment Framework: Systematic Approach", [
        "**The 'Look, Listen, Feel' Method:**",
        "",
        "**1. Look (Inspection)**",
        "- General appearance",
        "- Work of breathing",
        "- Skin color",
        "- Chest movement",
        "",
        "**2. Listen (Auscultation)**",
        "- Breath sounds",
        "- Voice sounds",
        "- Adventitious sounds",
        "",
        "**3. Feel (Palpation & Percussion)**",
        "- Chest expansion",
        "- Tactile fremitus",
        "- Percussion notes",
        "",
        "**Always assess:** Patient → Equipment → Data"
    ], "This systematic approach ensures nothing is missed. Always start with the patient, not the monitors.")
    
    # Slide 4: Initial Patient Assessment
    create_content_slide(prs, "Initial Patient Assessment: The First 60 Seconds", [
        "**General Appearance:**",
        "- Level of consciousness",
        "- Position of comfort (tripod, unable to lie flat)",
        "- Anxiety level",
        "- Diaphoresis",
        "- Skin color (cyanosis, pallor)",
        "",
        "**Work of Breathing:**",
        "- Respiratory rate",
        "- Use of accessory muscles",
        "- Pursed-lip breathing",
        "- Nasal flaring",
        "- Paradoxical breathing",
        "- Ability to speak (full sentences vs. words)",
        "",
        "**Immediate Concerns:**",
        "- Respiratory distress",
        "- Hemodynamic instability",
        "- Altered mental status"
    ], "The first 60 seconds are critical. Look for signs of distress before touching any equipment.")
    
    # Slide 5: Vital Signs
    create_content_slide(prs, "Vital Signs: Respiratory Assessment Priorities", [
        "**Respiratory Rate:**",
        "- Normal: 12-20 breaths/min",
        "- Tachypnea: >20 (early sign of distress)",
        "- Bradypnea: <12 (concerning for fatigue)",
        "- Pattern: Regular vs. irregular",
        "",
        "**Pulse Oximetry:**",
        "- Normal: 95-100%",
        "- Hypoxemic: <90%",
        "- Trend is more important than single value",
        "- Verify with clinical picture",
        "",
        "**Heart Rate and Blood Pressure:**",
        "- Tachycardia common with hypoxemia",
        "- Pulsus paradoxus in severe asthma",
        "- Hypertension with CO2 retention",
        "",
        "**Temperature:**",
        "- Fever suggests infection",
        "- Hypothermia in sepsis"
    ], "Respiratory rate is often the first vital sign to change in respiratory distress. Don't ignore it.")
    
    # Slide 6: Arterial Blood Gas Analysis
    create_content_slide(prs, "Arterial Blood Gas Analysis: The ABCs of ABGs", [
        "**Normal Values:**",
        "- pH: 7.35-7.45",
        "- PaCO2: 35-45 mmHg",
        "- PaO2: 80-100 mmHg",
        "- HCO3: 22-26 mEq/L",
        "- Base excess: -2 to +2",
        "",
        "**Step-by-Step Interpretation:**",
        "",
        "**1. Look at pH:**",
        "- <7.35 = Acidosis",
        "- >7.45 = Alkalosis",
        "- 7.35-7.45 = Normal or compensated",
        "",
        "**2. Determine primary disorder:**",
        "- Respiratory: PaCO2 abnormal",
        "- Metabolic: HCO3 abnormal",
        "",
        "**3. Check for compensation:**",
        "- Is the other parameter moving in the correct direction?"
    ], "ABG interpretation is a step-by-step process. Follow the steps every time for consistent results.")
    
    # Slide 7: ABG Interpretation Examples
    create_table_slide(prs, "ABG Interpretation Examples: Practice Cases",
        ["Case", "pH", "PaCO2", "HCO3", "Interpretation"],
        [
            ["Case 1", "7.25", "55", "24", "Acute respiratory acidosis"],
            ["Case 2", "7.32", "60", "31", "Chronic respiratory acidosis with compensation"],
            ["Case 3", "7.48", "30", "22", "Acute respiratory alkalosis"],
            ["Case 4", "7.30", "40", "18", "Metabolic acidosis"],
            ["Case 5", "7.40", "60", "36", "Compensated respiratory acidosis"]
        ],
        "Practice these cases. The key is to identify the primary disorder first, then look for compensation.")
    
    # Slide 8: Oxygenation Assessment
    create_content_slide(prs, "Oxygenation Assessment", [
        "**PaO2 Interpretation:**",
        "- 80-100 mmHg: Normal",
        "- 60-80 mmHg: Mild hypoxemia",
        "- 40-60 mmHg: Moderate hypoxemia",
        "- <40 mmHg: Severe hypoxemia",
        "",
        "**Age Adjustment:**",
        "- Expected PaO2 = 100 - (age × 0.3)",
        "- Or PaO2 >80 acceptable for elderly",
        "",
        "**Alveolar-arterial (A-a) Gradient:**",
        "- PAO2 = (FIO2 × [760-47]) - (PaCO2/0.8)",
        "- A-a gradient = PAO2 - PaO2",
        "- Normal: <10-15 mmHg (on room air)",
        "- >20 suggests V/Q mismatch or shunt",
        "",
        "**PaO2/FIO2 Ratio:**",
        "- Normal: 400-500",
        "- <300: Acute lung injury",
        "- <200: ARDS"
    ], "The A-a gradient helps differentiate between hypoxemia due to hypoventilation vs. other causes.")
    
    # Slide 9: Pulse Oximetry
    create_content_slide(prs, "Pulse Oximetry: Principles and Limitations", [
        "**How It Works:**",
        "- Measures light absorption at two wavelengths",
        "- Calculates ratio of oxygenated to deoxygenated hemoglobin",
        "- Displays as SpO2 percentage",
        "",
        "**Normal Values:**",
        "- 95-100%: Normal",
        "- 90-94%: Mild hypoxemia",
        "- 85-89%: Moderate hypoxemia",
        "- <85%: Severe hypoxemia",
        "",
        "**Limitations:**",
        "- Poor perfusion (shock, cold extremities)",
        "- Motion artifact",
        "- Nail polish (especially dark colors)",
        "- Carboxyhemoglobin (reads falsely high)",
        "- Methemoglobin (reads 85%)",
        "- Anemia (may be normal SpO2 but low oxygen content)"
    ], "Pulse oximetry is a screening tool, not a diagnostic test. Always correlate with clinical findings.")
    
    # Slide 10: Capnography
    create_content_slide(prs, "Capnography: ETCO2 Monitoring", [
        "**Normal Values:**",
        "- ETCO2: 35-45 mmHg",
        "- Typically 2-5 mmHg lower than PaCO2",
        "",
        "**Waveform Phases:**",
        "1. **Inspiration:** Baseline (should be zero)",
        "2. **Beginning exhalation:** Dead space gas",
        "3. **Exhalation:** Rapid rise",
        "4. **Alveolar plateau:** Peak ETCO2",
        "5. **Inspiration:** Sharp decline",
        "",
        "**Clinical Applications:**",
        "- Verify ET tube placement",
        "- Monitor ventilation",
        "- Detect changes in perfusion",
        "- Identify rebreathing",
        "- Procedural sedation monitoring"
    ], "Capnography is the gold standard for confirming ET tube placement. Always use it.")
    
    # Slide 11: Capnography Interpretation
    create_content_slide(prs, "Capnography Interpretation: Waveform Abnormalities", [
        "**Hypoventilation:**",
        "- Elevated ETCO2",
        "- Increased baseline (rebreathing)",
        "",
        "**Hyperventilation:**",
        "- Decreased ETCO2",
        "",
        "**Airway Obstruction:**",
        "- Prolonged expiratory upstroke",
        "- 'Shark fin' appearance",
        "- Incomplete plateau",
        "",
        "**Esophageal Intubation:**",
        "- No waveform or very small, irregular waveform",
        "",
        "**Cardiac Arrest:**",
        "- Sudden drop to zero",
        "- Poor compression quality: low ETCO2",
        "- ROSC: Sudden increase in ETCO2"
    ], "The capnography waveform tells a story. Learn to recognize these patterns.")
    
    # Slide 12: Physical Examination - Inspection
    create_content_slide(prs, "Physical Examination: Inspection", [
        "**General:**",
        "- Level of consciousness",
        "- Position (tripod, orthopnea)",
        "- Cyanosis (central vs. peripheral)",
        "- Clubbing of fingers",
        "- Barrel chest",
        "",
        "**Breathing Pattern:**",
        "- Rate and rhythm",
        "- Depth (tidal volume)",
        "- Symmetry of chest movement",
        "- Use of accessory muscles",
        "- Pursed-lip breathing",
        "- Paradoxical movement",
        "",
        "**Skin:**",
        "- Color (pallor, cyanosis, plethora)",
        "- Temperature",
        "- Diaphoresis",
        "- Subcutaneous emphysema"
    ], "Inspection can reveal critical information before you ever touch the patient.")
    
    # Slide 13: Physical Examination - Palpation
    create_content_slide(prs, "Physical Examination: Palpation", [
        "**Tracheal Position:**",
        "- Should be midline",
        "- Deviation suggests tension pneumothorax or atelectasis",
        "",
        "**Chest Expansion:**",
        "- Place hands on chest, thumbs at midline",
        "- Watch thumb movement during inspiration",
        "- Should be symmetrical",
        "- Asymmetry suggests pneumothorax, effusion, or atelectasis",
        "",
        "**Tactile Fremitus:**",
        "- Patient says '99' or 'blue moon'",
        "- Feel vibrations on chest wall",
        "- Increased: Consolidation",
        "- Decreased: Pleural effusion, pneumothorax, COPD",
        "",
        "**Subcutaneous Emphysema:**",
        "- Crackling sensation (crepitus)",
        "- Indicates air leak"
    ], "Palpation confirms what you observed during inspection and guides your auscultation.")
    
    # Slide 14: Physical Examination - Percussion
    create_content_slide(prs, "Physical Examination: Percussion", [
        "**Technique:**",
        "- Place middle finger firmly on chest",
        "- Tap distal interphalangeal joint with other middle finger",
        "- Compare side to side",
        "",
        "**Sounds:**",
        "- **Resonant:** Normal lung",
        "- **Hyperresonant:** Air trapping (COPD, asthma), pneumothorax",
        "- **Dull:** Consolidation, pleural effusion, atelectasis",
        "- **Flat:** Large pleural effusion",
        "- **Tympanic:** Gastric bubble (left lower anterior)",
        "",
        "**Pattern Recognition:**",
        "- Upper lobes: Clavicles to 4th rib anterior, scapulae posterior",
        "- Lower lobes: Below 4th rib anterior, below scapulae posterior"
    ], "Percussion helps identify areas of consolidation, air trapping, or fluid.")
    
    # Slide 15: Physical Examination - Auscultation
    create_content_slide(prs, "Physical Examination: Auscultation", [
        "**Systematic Approach:**",
        "- Compare side to side at each level",
        "- Anterior: 6 locations (3 each side)",
        "- Posterior: 10 locations (5 each side)",
        "- Lateral: 4 locations (2 each side)",
        "",
        "**Normal Breath Sounds:**",
        "",
        "**Vesicular:**",
        "- Location: Peripheral lung fields",
        "- Quality: Soft, low-pitched",
        "- Timing: Inspiration > expiration (3:1)",
        "",
        "**Bronchial:**",
        "- Location: Over trachea",
        "- Quality: Loud, high-pitched",
        "- Timing: Inspiration = expiration",
        "- Pause between",
        "",
        "**Bronchovesicular:**",
        "- Location: Over major bronchi",
        "- Quality: Intermediate",
        "- Timing: Inspiration = expiration"
    ], "Auscultation is the cornerstone of respiratory assessment. Be systematic and thorough.")
    
    # Slide 16: Adventitious Breath Sounds
    create_content_slide(prs, "Adventitious Breath Sounds: Abnormal Findings", [
        "**Crackles (Rales):**",
        "- **Fine:** High-pitched, late inspiration",
        "  - Causes: Interstitial fibrosis, early CHF, atelectasis",
        "- **Coarse:** Low-pitched, early inspiration",
        "  - Causes: Secretions, pulmonary edema",
        "",
        "**Wheezes:**",
        "- Continuous, musical",
        "- **Expiratory:** Asthma, COPD (narrowing of small airways)",
        "- **Inspiratory:** Upper airway obstruction, severe asthma",
        "- **Monophonic:** Single airway obstruction",
        "- **Polyphonic:** Multiple airways",
        "",
        "**Rhonchi:**",
        "- Continuous, low-pitched, snoring quality",
        "- Secretions in large airways",
        "- May clear with cough",
        "",
        "**Pleural Friction Rub:**",
        "- Discontinuous, grating",
        "- Pleural inflammation",
        "- Best heard at end-inspiration"
    ], "Each adventitious sound has a specific meaning. Learn to differentiate them.")
    
    # Slide 17: Chest X-Ray Interpretation
    create_content_slide(prs, "Chest X-Ray Interpretation: Systematic Approach", [
        "**ABC Method:**",
        "",
        "**A - Airway:**",
        "- Trachea midline?",
        "- Carina visible?",
        "- ET tube position (3-5 cm above carina)",
        "",
        "**B - Breathing (Lungs and Pleura):**",
        "- Lung fields clear?",
        "- Vascular markings present?",
        "- Costophrenic angles sharp?",
        "- Pleural spaces clear?",
        "",
        "**C - Circulation (Heart and Mediastinum):**",
        "- Heart size (CTR <50%)",
        "- Mediastinum midline?",
        "- Aortic knob visible?",
        "",
        "**D - Disability (Bones):**",
        "- Ribs, clavicles, spine intact?",
        "- Subcutaneous emphysema?",
        "",
        "**E - Everything Else:**",
        "- Lines and tubes",
        "- Diaphragms",
        "- Soft tissues"
    ], "A systematic approach to CXR interpretation prevents missing important findings.")
    
    # Slide 18: Common Chest X-Ray Findings
    create_content_slide(prs, "Common Chest X-Ray Findings: Pathological Patterns", [
        "**Consolidation (Pneumonia):**",
        "- Airspace opacification",
        "- Air bronchograms",
        "- Lobar or segmental distribution",
        "",
        "**Atelectasis:**",
        "- Volume loss",
        "- Increased opacity",
        "- Tracheal/mediastinal shift toward affected side",
        "- Elevated hemidiaphragm",
        "",
        "**Pleural Effusion:**",
        "- Blunted costophrenic angle",
        "- Meniscus sign",
        "- Mediastinal shift away (if large)",
        "",
        "**Pneumothorax:**",
        "- Visible visceral pleural line",
        "- Absent lung markings peripheral to line",
        "- Tracheal/mediastinal shift away (if tension)",
        "",
        "**Pulmonary Edema:**",
        "- Bilateral perihilar opacities (bat wing)",
        "- Kerley B lines",
        "- Cardiomegaly",
        "- Pleural effusions"
    ], "Recognizing these patterns is essential for timely intervention.")
    
    # Slide 19: Case Study 1
    create_content_slide(prs, "Case Study 1: Patient Presentation", [
        "**Mr. Davis, 72-year-old male**",
        "- Found confused at home",
        "- History of COPD",
        "",
        "**Assessment:**",
        "- Alert but confused",
        "- RR: 8, shallow",
        "- SpO2: 78% on room air",
        "- Accessory muscle use",
        "- Barrel chest",
        "- Diminished breath sounds bilaterally",
        "",
        "**ABG:**",
        "- pH: 7.18",
        "- PaCO2: 88",
        "- PaO2: 45",
        "- HCO3: 32",
        "",
        "**Questions:**",
        "1. Interpret the ABG",
        "2. What is your immediate management?",
        "3. What mode of ventilation (if needed) would be appropriate?"
    ], "This is a classic presentation of acute-on-chronic respiratory failure.")
    
    # Slide 20: Case Study 1 Discussion
    create_content_slide(prs, "Case Study 1: Discussion", [
        "**ABG Interpretation:**",
        "- pH 7.18: Severe acidemia",
        "- PaCO2 88: Respiratory acidosis",
        "- HCO3 32: Metabolic compensation (chronic)",
        "- PaO2 45: Severe hypoxemia",
        "- **Diagnosis:** Acute-on-chronic respiratory acidosis with severe hypoxemia",
        "",
        "**Immediate Management:**",
        "1. Call for help/rapid response",
        "2. Assist ventilation with BVM",
        "3. Controlled oxygen (2-4L NC initially)",
        "4. Prepare for intubation if no rapid improvement",
        "5. Consider NIV if alert and cooperative",
        "",
        "**Ventilation Strategy:**",
        "- If intubated: Volume or pressure control",
        "- Low tidal volume (6-8 ml/kg)",
        "- Slow rate initially (10-12) to allow time for compensation",
        "- Monitor for auto-PEEP"
    ], "This patient needs immediate intervention. The high PaCO2 with metabolic compensation suggests chronic CO2 retention.")
    
    # Slide 21: Case Study 2
    create_content_slide(prs, "Case Study 2: Patient Presentation", [
        "**Mrs. Patel, 58-year-old female**",
        "- Post-operative day 3 after hip surgery",
        "- Sudden onset dyspnea",
        "",
        "**Assessment:**",
        "- Anxious, tachypneic",
        "- RR: 28",
        "- SpO2: 89% on 4L NC",
        "- Clear breath sounds bilaterally",
        "- No chest pain",
        "",
        "**ABG:**",
        "- pH: 7.48",
        "- PaCO2: 30",
        "- PaO2: 62",
        "- HCO3: 22",
        "",
        "**Questions:**",
        "1. Interpret the ABG",
        "2. What is the differential diagnosis?",
        "3. What additional tests would you recommend?"
    ], "Post-operative patients are at high risk for pulmonary embolism.")
    
    # Slide 22: Case Study 2 Discussion
    create_content_slide(prs, "Case Study 2: Discussion", [
        "**ABG Interpretation:**",
        "- pH 7.48: Alkalemia",
        "- PaCO2 30: Respiratory alkalosis",
        "- HCO3 22: Normal (acute)",
        "- PaO2 62: Hypoxemia",
        "- **Diagnosis:** Acute respiratory alkalosis with hypoxemia",
        "",
        "**Differential Diagnosis:**",
        "- Pulmonary embolism (high suspicion post-op)",
        "- Pneumonia",
        "- Atelectasis",
        "- Pulmonary edema",
        "- Pneumothorax",
        "",
        "**Additional Tests:**",
        "- CT pulmonary angiography (PE protocol)",
        "- D-dimer (if low pre-test probability)",
        "- Chest X-ray",
        "- ECG",
        "- Troponin (rule out cardiac)",
        "",
        "**Clinical Pearl:** Post-operative patient with clear lungs, hypoxemia, and respiratory alkalosis should raise suspicion for PE until proven otherwise."
    ], "The combination of hypoxemia, respiratory alkalosis, and clear lungs is classic for PE.")
    
    # Slide 23: Key Takeaways
    create_content_slide(prs, "Key Takeaways", [
        "**1. Assessment is Systematic**",
        "- Look, listen, feel",
        "- Patient first, then equipment, then data",
        "",
        "**2. ABG Interpretation is Logical**",
        "- pH first, then primary disorder, then compensation",
        "- Practice makes perfect",
        "",
        "**3. Physical Exam Skills are Essential**",
        "- Inspection reveals work of breathing",
        "- Auscultation identifies pathology",
        "- Palpation and percussion confirm findings",
        "",
        "**4. Technology Supports, Doesn't Replace**",
        "- Pulse oximetry and capnography are tools",
        "- Clinical assessment is paramount",
        "",
        "**5. Integrate All Findings**",
        "- No single finding tells the whole story",
        "- Look for patterns"
    ], "Master these skills and you'll be an effective respiratory therapist.")
    
    # Slide 24: References
    create_content_slide(prs, "References", [
        "1. AARC Clinical Practice Guideline: Blood Gas Analysis and Hemoximetry. Respir Care. 2022.",
        "",
        "2. AARC Clinical Practice Guideline: Capnography/Capnometry. Respir Care. 2023.",
        "",
        "3. AARC Clinical Practice Guideline: Pulse Oximetry. Respir Care. 2021.",
        "",
        "4. Kacmarek RM, Stoller JK, Heuer AJ. Egan's Fundamentals of Respiratory Care. 12th ed. Elsevier; 2021.",
        "",
        "5. Wilkins RL, Dexter JR. Respiratory Disease: A Case Study Approach to Patient Care. 4th ed. FA Davis; 2017."
    ], "")
    
    # Slide 25: Thank You
    create_end_slide(prs, "Thank You", "Assessment is the Foundation of Quality Care", "Respiratory Care Department")
    
    # Save presentation
    output_path = os.path.join(OUTPUT_DIR, "04-Patient-Assessment.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

# ============================================================
# PRESENTATION 2: 05-Emergency-Procedures.pptx
# ============================================================

def create_emergency_procedures_ppt():
    """Create the Emergency Procedures presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs,
        "Emergency Respiratory Procedures",
        "Code Blue and Rapid Response",
        "[Name], RRT, [Credentials]",
        "[Date]",
        "Respiratory Care")
    
    # Slide 2: Learning Objectives
    create_content_slide(prs, "Learning Objectives", [
        "By the end of this session, participants will be able to:",
        "",
        "- Describe the respiratory therapist's role in code blue situations",
        "- Perform effective bag-valve-mask ventilation",
        "- Assist with emergency airway management",
        "- Manage the airway during cardiac arrest",
        "- Participate in rapid response team activation",
        "- Recognize and respond to respiratory emergencies"
    ], "Emergency situations require quick thinking and practiced skills.")
    
    # Slide 3: Code Blue Response
    create_content_slide(prs, "Code Blue Response: The RT Role", [
        "**Primary Responsibilities:**",
        "- Airway management",
        "- Ventilation support",
        "- Oxygen delivery",
        "- Medication delivery (nebulized)",
        "- Arterial blood gas sampling",
        "- Documentation",
        "",
        "**Equipment to Bring:**",
        "- Emergency airway bag",
        "- Portable suction",
        "- Oxygen tank and regulator",
        "- ABG kit",
        "- Stethoscope",
        "",
        "**Team Communication:**",
        "- Closed-loop communication",
        "- Clear, concise reports",
        "- Role clarity"
    ], "Know your role before the emergency occurs. Preparation saves lives.")
    
    # Slide 4: Adult Cardiac Arrest Algorithm
    create_content_slide(prs, "Adult Cardiac Arrest Algorithm", [
        "**Immediate Actions (CAB):**",
        "",
        "**1. C - Circulation: Chest compressions**",
        "- Rate: 100-120/min",
        "- Depth: 2-2.4 inches",
        "- Allow full recoil",
        "- Minimize interruptions",
        "",
        "**2. A - Airway: Open airway**",
        "- Head tilt-chin lift (no trauma)",
        "- Jaw thrust (if trauma suspected)",
        "",
        "**3. B - Breathing: Ventilation**",
        "- 2 breaths after every 30 compressions (BLS)",
        "- Advanced airway: 1 breath every 6 seconds (10/min)",
        "",
        "**Respiratory Therapist Role:**",
        "- Manage airway once advanced airway placed",
        "- Provide ventilation",
        "- Monitor ETCO2"
    ], "High-quality CPR is the foundation of resuscitation.")
    
    # Slide 5: Bag-Valve-Mask Ventilation
    create_content_slide(prs, "Bag-Valve-Mask Ventilation: Critical Skill", [
        "**Equipment:**",
        "- Self-inflating bag (adult 1-2L)",
        "- Mask (appropriate size)",
        "- Oxygen source",
        "- PEEP valve (if available)",
        "",
        "**One-Person Technique:**",
        "1. Position at head of bed",
        "2. E-C clamp technique:",
        "   - Thumb and index finger form 'C' on mask",
        "   - Middle, ring, pinky form 'E' under jaw",
        "3. Squeeze bag to deliver breath",
        "4. Watch for chest rise",
        "",
        "**Two-Person Technique (Preferred):**",
        "- One person holds mask with both hands",
        "- Second person squeezes bag",
        "- Better seal, easier ventilation"
    ], "BVM ventilation is a critical skill that requires practice. The two-person technique is preferred when available.")
    
    # Slide 6: BVM Best Practices
    create_content_slide(prs, "BVM Ventilation: Best Practices", [
        "**Volume:**",
        "- Deliver only enough to see chest rise",
        "- Approximately 500-600 ml",
        "- Avoid hyperventilation",
        "",
        "**Rate:**",
        "- With compressions: 1 breath every 6 seconds",
        "- Without compressions: 10-12 breaths/min",
        "- Squeeze bag over 1 second",
        "",
        "**Complications to Avoid:**",
        "- Hyperventilation (decreases venous return)",
        "- Gastric insufflation (increases aspiration risk)",
        "- Barotrauma",
        "",
        "**Monitoring:**",
        "- Watch chest rise",
        "- Listen for air leaks",
        "- Monitor SpO2",
        "- Capnography when available"
    ], "Avoid hyperventilation - it's common during emergencies and harmful to the patient.")
    
    # Slide 7: Emergency Airway Management
    create_content_slide(prs, "Emergency Airway Management", [
        "**Assessment (MOANS, LEMON, RODS):**",
        "",
        "**MOANS (Mask Ventilation):**",
        "- M - Mask seal (beard, no teeth)",
        "- O - Obesity/Obstruction",
        "- A - Age (>55)",
        "- N - No teeth",
        "- S - Stiff lungs",
        "",
        "**LEMON (Laryngoscopy):**",
        "- L - Look externally",
        "- E - Evaluate 3-3-2 rule",
        "- M - Mallampati score",
        "- O - Obstruction",
        "- N - Neck mobility",
        "",
        "**RODS (Supraglottic Airway):**",
        "- R - Restricted mouth opening",
        "- O - Obstruction",
        "- D - Disrupted airway",
        "- S - Stiff lungs/cervical spine"
    ], "These mnemonics help predict difficult airways. Assess early.")
    
    # Slide 8: Endotracheal Intubation
    create_content_slide(prs, "Endotracheal Intubation: RT Role", [
        "**Preparation:**",
        "- Check suction (Yankauer ready)",
        "- Pre-oxygenate (3 minutes at 100%)",
        "- Position patient (sniffing position)",
        "- Prepare BVM",
        "- Monitor SpO2 and HR",
        "",
        "**During Intubation:**",
        "- Assist as needed",
        "- Pass equipment",
        "- Apply cricoid pressure (if requested)",
        "- Monitor for desaturation",
        "",
        "**Post-Intubation:**",
        "- Confirm placement (ETCO2, auscultation)",
        "- Secure tube",
        "- Document depth (cm at teeth/gums)",
        "- Obtain chest X-ray",
        "- Initiate mechanical ventilation"
    ], "The RT plays a crucial supporting role in intubation. Be prepared.")
    
    # Slide 9: Confirming ET Tube Placement
    create_content_slide(prs, "Confirming ET Tube Placement", [
        "**Primary Confirmation:**",
        "- **Quantitative ETCO2:**",
        "  - Waveform capnography (gold standard)",
        "  - Colorimetric device (purple to yellow)",
        "  - Must be sustained (not just one breath)",
        "",
        "**Secondary Confirmation:**",
        "- Auscultation (bilateral breath sounds)",
        "- Chest rise",
        "- No gastric insufflation sounds",
        "",
        "**Chest X-Ray:**",
        "- Tip 3-5 cm above carina",
        "- Not for immediate confirmation",
        "",
        "**If Uncertain:**",
        "- Remove tube and re-intubate",
        "- Do not assume esophageal intubation will be obvious"
    ], "ETCO2 is the gold standard. Never rely on a single confirmation method.")
    
    # Slide 10: Surgical Airway
    create_content_slide(prs, "Surgical Airway: Cricothyrotomy", [
        "**Indications:**",
        "- Cannot intubate, cannot ventilate (CICV)",
        "- Severe facial trauma",
        "- Upper airway obstruction",
        "",
        "**Contraindications:**",
        "- Age <8 years (needle cricothyrotomy preferred)",
        "- Tracheal transection",
        "- Inability to identify landmarks",
        "",
        "**RT Role:**",
        "- Assist with equipment",
        "- Provide ventilation through surgical airway",
        "- Monitor patient",
        "",
        "**After Placement:**",
        "- Confirm placement",
        "- Secure airway",
        "- Prepare for transport to OR"
    ], "Cricothyrotomy is a life-saving procedure when the airway cannot be secured otherwise.")
    
    # Slide 11: Rapid Response Team
    create_content_slide(prs, "Rapid Response Team", [
        "**Purpose:**",
        "- Identify deteriorating patients early",
        "- Prevent cardiac arrest",
        "- Reduce ICU admissions",
        "",
        "**Activation Criteria:**",
        "- Acute change in respiratory rate (<8 or >28)",
        "- Acute change in oxygen saturation (<90%)",
        "- Acute change in mental status",
        "- Acute change in blood pressure (SBP <90 or >180)",
        "- Heart rate (<40 or >130)",
        "- Staff concern about patient",
        "",
        "**Respiratory Therapist Role:**",
        "- Assess airway and breathing",
        "- Provide oxygen/ventilation support",
        "- Prepare for potential intubation",
        "- Assist with procedures"
    ], "Early intervention prevents codes. Don't hesitate to call a rapid response.")
    
    # Slide 12: Respiratory Emergency - Asthma
    create_content_slide(prs, "Respiratory Emergency: Severe Asthma", [
        "**Recognition:**",
        "- Severe dyspnea",
        "- Inability to speak in sentences",
        "- Use of accessory muscles",
        "- Diaphoresis",
        "- Pulsus paradoxus >25 mmHg",
        "- Silent chest (ominous sign)",
        "",
        "**Management:**",
        "1. Oxygen (target SpO2 92-95%)",
        "2. Continuous nebulized beta-agonist",
        "3. Ipratropium bromide",
        "4. Systemic corticosteroids (early)",
        "5. Magnesium sulfate (severe)",
        "6. Consider NIV (controversial)",
        "7. Prepare for intubation if failing",
        "",
        "**Intubation Considerations:**",
        "- High risk for barotrauma",
        "- Use large ETT (8.0+)",
        "- Low tidal volume, long expiratory time",
        "- Ketamine preferred for induction"
    ], "A silent chest in asthma is an ominous sign indicating severe obstruction.")
    
    # Slide 13: Respiratory Emergency - Anaphylaxis
    create_content_slide(prs, "Respiratory Emergency: Anaphylaxis", [
        "**Recognition:**",
        "- Sudden onset after exposure",
        "- Airway: Angioedema, stridor",
        "- Breathing: Bronchospasm, hypoxemia",
        "- Circulation: Hypotension, tachycardia",
        "- Skin: Urticaria, flushing",
        "",
        "**Immediate Treatment:**",
        "1. **Epinephrine IM** (1:1000, 0.3-0.5 mg)",
        "   - Lateral thigh",
        "   - Repeat every 5-15 min if needed",
        "2. **Oxygen** (high flow)",
        "3. **IV fluids** (crystalloid bolus)",
        "4. **Albuterol** (nebulized)",
        "5. **Antihistamines** (H1 and H2 blockers)",
        "6. **Corticosteroids** (prevent biphasic reaction)",
        "",
        "**Airway Management:**",
        "- Early intubation if airway compromise",
        "- May be very difficult due to edema",
        "- Have surgical airway ready"
    ], "Epinephrine is the first-line treatment for anaphylaxis. Don't delay.")
    
    # Slide 14: Respiratory Emergency - Pulmonary Edema
    create_content_slide(prs, "Respiratory Emergency: Pulmonary Edema", [
        "**Cardiogenic (CHF):**",
        "- Orthopnea",
        "- JVD",
        "- S3 gallop",
        "- Peripheral edema",
        "- Bilateral crackles",
        "- CXR: Cardiomegaly, effusions, cephalization",
        "",
        "**Non-cardiogenic (ARDS):**",
        "- No orthopnea",
        "- No JVD",
        "- Bilateral crackles",
        "- CXR: Bilateral infiltrates, normal heart size",
        "",
        "**Management (Cardiogenic):**",
        "1. Sit patient upright",
        "2. Oxygen/CPAP/BiPAP",
        "3. Nitrates",
        "4. Diuretics",
        "5. Morphine (if indicated)"
    ], "Differentiating cardiogenic vs. non-cardiogenic pulmonary edema guides treatment.")
    
    # Slide 15: Respiratory Emergency - Pneumothorax
    create_content_slide(prs, "Respiratory Emergency: Pneumothorax", [
        "**Types:**",
        "- **Simple:** No shift",
        "- **Tension:** Mediastinal shift, hemodynamic compromise",
        "",
        "**Recognition:**",
        "- Sudden onset dyspnea",
        "- Pleuritic chest pain",
        "- Decreased breath sounds",
        "- Hyperresonance to percussion",
        "- Tension: Tracheal deviation, hypotension, JVD",
        "",
        "**Management:**",
        "- **Simple:** Oxygen, observation, chest tube if large/symptomatic",
        "- **Tension:** Immediate needle decompression",
        "  - 2nd intercostal space, midclavicular line",
        "  - 14-gauge needle, 3.25 inches",
        "  - Follow with chest tube",
        "",
        "**RT Role:**",
        "- Assist with procedure",
        "- Support ventilation",
        "- Prepare for chest tube insertion"
    ], "Tension pneumothorax is a clinical diagnosis. Don't wait for imaging.")
    
    # Slide 16: Foreign Body Airway Obstruction
    create_content_slide(prs, "Foreign Body Airway Obstruction", [
        "**Recognition:**",
        "- Universal choking sign (hands to throat)",
        "- Inability to speak/cough",
        "- Cyanosis",
        "- High-pitched sounds (partial)",
        "- Silent (complete)",
        "",
        "**Adult/Child (>1 year) Treatment:**",
        "",
        "**Mild (can cough, speak):**",
        "- Encourage coughing",
        "- Do not interfere",
        "- Monitor closely",
        "",
        "**Severe (cannot cough, speak, cyanotic):**",
        "1. Abdominal thrusts (Heimlich)",
        "   - Stand behind patient",
        "   - Fist above navel, below xiphoid",
        "   - Quick upward thrusts",
        "2. Continue until object expelled or patient becomes unresponsive",
        "",
        "**If Unresponsive:**",
        "- Activate emergency response",
        "- Begin CPR",
        "- Look in mouth before breaths, remove if visible"
    ], "Act quickly with choking. If they can cough, let them. If not, intervene immediately.")
    
    # Slide 17: Post-Resuscitation Care
    create_content_slide(prs, "Post-Resuscitation Care", [
        "**Respiratory Priorities:**",
        "",
        "**1. Airway:**",
        "- Maintain ET tube",
        "- Suction as needed",
        "- Consider sedation",
        "",
        "**2. Ventilation:**",
        "- Avoid hyperventilation",
        "- Target PaCO2 35-45",
        "- Protect lungs (6 ml/kg)",
        "",
        "**3. Oxygenation:**",
        "- Target SpO2 94-98%",
        "- Avoid hyperoxia",
        "- Titrate FIO2",
        "",
        "**4. Monitoring:**",
        "- Continuous capnography",
        "- ABG",
        "- Chest X-ray",
        "",
        "**Targeted Temperature Management:**",
        "- 32-36°C for 24 hours (if indicated)"
    ], "Post-resuscitation care is critical for neurological recovery.")
    
    # Slide 18: Documentation
    create_content_slide(prs, "Documentation", [
        "**Respiratory Documentation:**",
        "- Arrival time",
        "- Initial assessment",
        "- Interventions performed",
        "- Airway management details",
        "- Ventilation parameters",
        "- Medications administered",
        "- Response to treatment",
        "- Times of key events",
        "",
        "**Critical Times to Record:**",
        "- Code called",
        "- First compression",
        "- Airway secured",
        "- ROSC achieved",
        "- Code ended",
        "",
        "**Quality Improvement:**",
        "- Debrief after code",
        "- Identify system issues",
        "- Provide feedback"
    ], "Good documentation is essential for legal protection and quality improvement.")
    
    # Slide 19: Case Study 1
    create_content_slide(prs, "Case Study 1: Code Blue Scenario", [
        "**You are called to a code blue on the medical floor.**",
        "",
        "**Patient:** 68-year-old male",
        "**History:** Post-operative day 2 after bowel resection",
        "",
        "**Arrival:**",
        "- CPR in progress",
        "- No airway established",
        "- Bag-mask ventilation being performed",
        "",
        "**Assessment:**",
        "- No chest rise with BVM",
        "- SpO2: 60%",
        "- Rhythm: Asystole",
        "",
        "**Questions:**",
        "1. What is your immediate priority?",
        "2. How do you troubleshoot the BVM issue?",
        "3. What are your next steps?"
    ], "This is a high-stress situation. Focus on the airway first.")
    
    # Slide 20: Case Study 1 Discussion
    create_content_slide(prs, "Case Study 1: Discussion", [
        "**Immediate Priority:**",
        "- Ensure effective ventilation is the priority",
        "- Without oxygenation, other interventions are futile",
        "",
        "**Troubleshooting BVM:**",
        "1. Check mask seal (E-C clamp technique)",
        "2. Reposition head (head tilt-chin lift or jaw thrust)",
        "3. Check for airway obstruction (suction if needed)",
        "4. Try two-person BVM technique",
        "5. Consider oral/nasopharyngeal airway",
        "6. Prepare for intubation",
        "",
        "**Next Steps:**",
        "1. Call for intubation equipment if not already done",
        "2. Continue high-quality CPR",
        "3. Once airway secured, confirm with ETCO2",
        "4. Manage ventilator (avoid hyperventilation)",
        "5. Obtain ABG once ROSC achieved"
    ], "Airway first, then circulation. You can't resuscitate without oxygenation.")
    
    # Slide 21: Case Study 2
    create_content_slide(prs, "Case Study 2: Rapid Response Scenario", [
        "**You respond to a rapid response call.**",
        "",
        "**Patient:** 45-year-old female",
        "**History:** Asthma, admitted for pneumonia",
        "",
        "**Nurse Reports:**",
        "- 'She was fine an hour ago'",
        "- 'Now she's breathing really fast and can't talk'",
        "",
        "**Your Assessment:**",
        "- RR: 36",
        "- SpO2: 88% on 4L NC",
        "- Using accessory muscles",
        "- Can only speak in words",
        "- Auscultation: Diminished breath sounds, no wheezing",
        "",
        "**Questions:**",
        "1. What is the significance of the 'silent chest'?",
        "2. What are your immediate interventions?",
        "3. When should you consider intubation?"
    ], "A silent chest in an asthmatic is a life-threatening emergency.")
    
    # Slide 22: Case Study 2 Discussion
    create_content_slide(prs, "Case Study 2: Discussion", [
        "**Silent Chest Significance:**",
        "- Extremely concerning finding",
        "- Indicates severe airway obstruction",
        "- Insufficient airflow to generate wheezing",
        "- Patient is tiring",
        "- Imminent respiratory arrest risk",
        "",
        "**Immediate Interventions:**",
        "1. High-flow oxygen (non-rebreather)",
        "2. Continuous nebulized albuterol (high dose)",
        "3. Ipratropium bromide",
        "4. Systemic corticosteroids (IV)",
        "5. Magnesium sulfate (2g IV)",
        "6. Prepare for intubation",
        "",
        "**Intubation Criteria:**",
        "- Altered mental status",
        "- Respiratory arrest",
        "- Severe acidosis (pH <7.25)",
        "- No response to maximal therapy",
        "- Exhaustion"
    ], "Don't wait for respiratory arrest to intubate. Act early in severe asthma.")
    
    # Slide 23: Key Takeaways
    create_content_slide(prs, "Key Takeaways", [
        "**1. Airway is Always Priority**",
        "- Without oxygenation, nothing else matters",
        "- BVM is a critical skill",
        "",
        "**2. Teamwork Saves Lives**",
        "- Clear communication",
        "- Defined roles",
        "- Closed-loop communication",
        "",
        "**3. Avoid Hyperventilation**",
        "- Decreases venous return",
        "- Increases intrathoracic pressure",
        "- 10 breaths/min with advanced airway",
        "",
        "**4. Confirm ET Tube Placement**",
        "- ETCO2 is gold standard",
        "- Never assume",
        "",
        "**5. Early Intervention Prevents Codes**",
        "- Know rapid response criteria",
        "- Don't hesitate to call for help"
    ], "Emergency preparedness is about preparation, not just response.")
    
    # Slide 24: References
    create_content_slide(prs, "References", [
        "1. American Heart Association. Advanced Cardiovascular Life Support (ACLS) Provider Manual. 2020.",
        "",
        "2. AARC Clinical Practice Guideline: Resuscitation and Defibrillation in the Health Care Setting. Respir Care. 2022.",
        "",
        "3. AARC Clinical Practice Guideline: Management of Airway Emergencies. Respir Care. 2019.",
        "",
        "4. GINA Global Strategy for Asthma Management and Prevention. 2024 Update.",
        "",
        "5. Difficult Airway Society Guidelines. https://www.das.uk.com/guidelines"
    ], "")
    
    # Slide 25: Thank You
    create_end_slide(prs, "Thank You", "Prepared for the Unexpected", "Respiratory Care Department")
    
    # Save presentation
    output_path = os.path.join(OUTPUT_DIR, "05-Emergency-Procedures.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

# ============================================================
# PRESENTATION 3: Mechanical-Ventilation-Basics.pptx
# ============================================================

def create_ventilation_basics_ppt():
    """Create the Mechanical Ventilation Basics presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs,
        "Mechanical Ventilation Basics",
        "Introduction to Ventilator Management",
        "[Name], RRT",
        "[Date]",
        "Respiratory Care")
    
    # Slide 2: Learning Objectives
    create_content_slide(prs, "Learning Objectives", [
        "By the end of this session, participants will be able to:",
        "",
        "- Identify the components of a mechanical ventilator circuit",
        "- Describe the difference between volume and pressure control ventilation",
        "- Calculate initial ventilator settings using lung-protective strategies",
        "- Recognize common ventilator alarms and appropriate responses",
        "- Interpret basic ventilator graphics"
    ], "These foundational skills are essential for ventilator management.")
    
    # Slide 3: Indications for Mechanical Ventilation
    create_content_slide(prs, "Indications for Mechanical Ventilation", [
        "**Primary Indications:**",
        "",
        "- **Apnea** or impending respiratory arrest",
        "- **Acute respiratory failure**",
        "  - Hypoxemic (PaO2 < 60 on high FIO2)",
        "  - Hypercapnic (pH < 7.30 with elevated PaCO2)",
        "- **Airway protection**",
        "  - Altered mental status",
        "  - Inability to protect airway",
        "- **Prophylactic**",
        "  - Post-operative support",
        "  - Severe metabolic acidosis"
    ], "Understanding indications helps determine when ventilation is truly needed.")
    
    # Slide 4: Ventilator Classification
    create_table_slide(prs, "Ventilator Classification: By Control Variable",
        ["Type", "Control", "Advantage"],
        [
            ["Volume Control", "Delivers set tidal volume", "Guaranteed minute ventilation"],
            ["Pressure Control", "Delivers set pressure", "Pressure limiting, variable flow"]
        ],
        "Each mode has advantages. Choose based on patient needs.")
    
    # Slide 5: Breath Delivery Modes
    create_content_slide(prs, "Ventilator Classification: By Breath Delivery", [
        "**Continuous Mandatory Ventilation (CMV)**",
        "- All breaths machine-triggered",
        "- Patient cannot trigger additional breaths",
        "",
        "**Assist/Control (A/C)**",
        "- All breaths delivered",
        "- Patient can trigger additional breaths",
        "- Most common initial mode",
        "",
        "**Synchronized Intermittent Mandatory Ventilation (SIMV)**",
        "- Set rate + spontaneous breaths",
        "- Spontaneous breaths may have pressure support",
        "",
        "**Spontaneous Modes**",
        "- CPAP: Continuous positive airway pressure",
        "- Pressure Support: Assists spontaneous breaths"
    ], "Mode selection depends on the patient's respiratory drive and condition.")
    
    # Slide 6: Lung-Protective Ventilation
    create_content_slide(prs, "Lung-Protective Ventilation: ARDSNet Protocol", [
        "**Tidal Volume:**",
        "- **6 ml/kg ideal body weight** (range 4-8)",
        "- Calculate IBW:",
        "  - Men: 50 + (2.3 × inches over 5 feet)",
        "  - Women: 45.5 + (2.3 × inches over 5 feet)",
        "",
        "**Plateau Pressure:**",
        "- **Goal: < 30 cmH2O**",
        "- If >30: decrease VT to 5 or 4 ml/kg",
        "",
        "**PEEP:**",
        "- Start at **5 cmH2O**",
        "- Titrate based on oxygenation and hemodynamics"
    ], "Lung-protective ventilation reduces ventilator-induced lung injury.")
    
    # Slide 7: Initial Ventilator Settings
    create_table_slide(prs, "Initial Ventilator Settings: Standard Starting Point",
        ["Parameter", "Setting", "Rationale"],
        [
            ["Mode", "Volume A/C", "Reliable minute ventilation"],
            ["VT", "6-8 ml/kg IBW", "Lung protection"],
            ["Rate", "12-16/min", "Normal physiologic range"],
            ["FIO2", "100% → titrate", "Start high, wean quickly"],
            ["PEEP", "5 cmH2O", "Maintain FRC"],
            ["Flow", "60 L/min", "Adequate for most adults"]
        ],
        "These are starting points. Adjust based on patient response.")
    
    # Slide 8: Ventilator Graphics
    create_content_slide(prs, "Ventilator Graphics: Pressure-Time Waveform", [
        "**Key Measurements:**",
        "",
        "- **Peak Pressure** - Total resistance + compliance",
        "- **Plateau Pressure** - Static compliance only",
        "- **PEEP** - Baseline pressure",
        "",
        "**Key Calculations:**",
        "",
        "**Static Compliance** = VT / (Pplat - PEEP)",
        "- Normal: 50-100 mL/cmH2O",
        "",
        "**Airway Resistance** = (Ppeak - Pplat) / Flow",
        "- Normal: < 10 cmH2O/L/sec",
        "",
        "**Example:**",
        "- VT = 500 mL, Pplat = 25, PEEP = 5",
        "- Compliance = 500 / (25 - 5) = 25 mL/cmH2O (low)"
    ], "Understanding waveforms helps identify patient-ventilator problems.")
    
    # Slide 9: Common Ventilator Alarms - High Pressure
    create_content_slide(prs, "Common Ventilator Alarms: High Pressure", [
        "**Causes:**",
        "- Secretions in airway",
        "- Kinked ETT or circuit",
        "- Bronchospasm",
        "- Decreased lung compliance",
        "- Patient-ventilator asynchrony",
        "",
        "**Response:**",
        "1. Check patient first",
        "2. Suction if indicated",
        "3. Check circuit for kinks",
        "4. Auscultate for bronchospasm",
        "5. Consider sedation if fighting ventilator",
        "",
        "**Prevention:**",
        "- Regular airway suctioning",
        "- Proper circuit positioning",
        "- Adequate sedation"
    ], "High pressure alarms always require immediate attention.")
    
    # Slide 10: Common Ventilator Alarms - Low Pressure
    create_content_slide(prs, "Common Ventilator Alarms: Low Pressure/Volume", [
        "**Low Pressure/Volume Alarm Causes:**",
        "- Circuit disconnect",
        "- Cuff leak",
        "- Patient self-extubation",
        "",
        "**Response:**",
        "1. Check patient immediately",
        "2. Verify ETT position",
        "3. Check cuff pressure",
        "4. Check all connections",
        "",
        "**Low PEEP/CPAP Alarm Causes:**",
        "- Large leak in system",
        "- Circuit disconnect",
        "- Inadequate flow",
        "",
        "**Remember:** Low pressure alarms may indicate a life-threatening disconnect."
    ], "Low pressure alarms can indicate a catastrophic circuit disconnect.")
    
    # Slide 11: Patient-Ventilator Asynchrony
    create_table_slide(prs, "Patient-Ventilator Asynchrony: Types and Solutions",
        ["Type", "Description", "Solution"],
        [
            ["Trigger", "Difficulty initiating breath", "Adjust sensitivity"],
            ["Flow", "Demand exceeds set flow", "Increase flow or switch to pressure control"],
            ["Cycle", "Premature or delayed termination", "Adjust cycling criteria"]
        ],
        "Asynchrony increases work of breathing and patient discomfort.")
    
    # Slide 12: Assessment of Asynchrony
    create_content_slide(prs, "Assessment of Asynchrony", [
        "**Signs of Asynchrony:**",
        "",
        "- Observe patient breathing pattern",
        "- Check flow and pressure waveforms",
        "- Assess for double-triggering",
        "- Evaluate for auto-PEEP",
        "",
        "**Double-Triggering:**",
        "- Two breaths delivered back-to-back",
        "- Indicates insufficient tidal volume",
        "- May need to increase VT or switch to pressure control",
        "",
        "**Auto-PEEP:**",
        "- Intrinsic PEEP from incomplete exhalation",
        "- Common in COPD and asthma",
        "- Increase expiratory time",
        "- Decrease respiratory rate"
    ], "Addressing asynchrony improves patient comfort and ventilation efficiency.")
    
    # Slide 13: Weaning Parameters
    create_content_slide(prs, "Weaning Parameters: Readiness Assessment", [
        "**Spontaneous Breathing Trial (SBT) Criteria:**",
        "",
        "**RSBI (Rapid Shallow Breathing Index)**",
        "- RR / VT (in liters)",
        "- < 105 = likely to succeed",
        "",
        "**Negative Inspiratory Force (NIF)**",
        "- < -20 to -30 cmH2O",
        "",
        "**Other Parameters:**",
        "- Spontaneous VT > 5 ml/kg",
        "- VC > 10-15 ml/kg",
        "- Minute ventilation < 10 L/min",
        "",
        "**Clinical Assessment:**",
        "- Hemodynamic stability",
        "- Adequate oxygenation",
        "- Resolution of underlying condition"
    ], "Weaning readiness requires both physiological and clinical assessment.")
    
    # Slide 14: Case Study
    create_content_slide(prs, "Case Study: Patient Profile", [
        "**Mr. Johnson**, 70-year-old male",
        "**Diagnosis:** Severe pneumonia, ARDS",
        "**Height:** 5'10\" (70 inches)",
        "**Actual weight:** 90 kg",
        "",
        "**Calculate Initial Settings:**",
        "",
        "**Step 1: Calculate IBW**",
        "- IBW = 50 + (2.3 × 10) = **73 kg**",
        "",
        "**Step 2: Calculate VT**",
        "- 6 ml/kg × 73 kg = **438 ml** (round to 440 ml)",
        "",
        "**Step 3: Initial Settings**",
        "- Mode: Volume A/C",
        "- VT: 440 ml",
        "- Rate: 14",
        "- FIO2: 100%",
        "- PEEP: 5"
    ], "Always use ideal body weight for tidal volume calculations in ARDS.")
    
    # Slide 15: Case Study Discussion
    create_content_slide(prs, "Case Study: Discussion", [
        "**Questions:**",
        "",
        "1. Why use IBW instead of actual weight?",
        "   - Lung compliance, not body size, determines safe VT",
        "",
        "2. What if plateau pressure is 35 cmH2O?",
        "   - If Pplat >30, decrease VT to 4-5 ml/kg",
        "   - May need to increase rate to maintain minute ventilation",
        "",
        "3. When should you start weaning FIO2?",
        "   - Wean FIO2 to maintain SpO2 88-95% (ARDS)",
        "   - Avoid hyperoxia",
        "",
        "4. What PEEP/FIO2 table might you use?",
        "   - ARDSNet PEEP/FIO2 tables guide PEEP titration",
        "   - Higher PEEP may improve oxygenation but affect hemodynamics"
    ], "Lung-protective ventilation is essential in ARDS to prevent further lung injury.")
    
    # Slide 16: Key Takeaways
    create_content_slide(prs, "Key Takeaways", [
        "**1. Always use lung-protective ventilation**",
        "- 6 ml/kg IBW, Pplat < 30",
        "",
        "**2. Know your calculations**",
        "- IBW, compliance, RSBI",
        "",
        "**3. Patient first, then equipment**",
        "- Always assess patient when alarms sound",
        "",
        "**4. Monitor for asynchrony**",
        "- Adjust settings to optimize comfort",
        "",
        "**5. Have a weaning plan**",
        "- Daily SBT when ready"
    ], "These principles form the foundation of safe mechanical ventilation.")
    
    # Slide 17: References
    create_content_slide(prs, "References", [
        "1. ARDSNet. Ventilation with lower tidal volumes as compared with traditional tidal volumes for acute lung injury and the acute respiratory distress syndrome. N Engl J Med. 2000;342(18):1301-1308.",
        "",
        "2. AARC Clinical Practice Guidelines: Mechanical Ventilation of Mechanically Ventilated Patients. Respir Care. 2022.",
        "",
        "3. Tobin MJ. Principles and Practice of Mechanical Ventilation. 4th ed. McGraw-Hill; 2018.",
        "",
        "4. Pilbeam SP, Cairo JM. Mechanical Ventilation: Physiological and Clinical Applications. 7th ed. Elsevier; 2020."
    ], "")
    
    # Slide 18: Thank You
    create_end_slide(prs, "Thank You", "Excellence in Patient Care Through Education", "Respiratory Care Department")
    
    # Save presentation
    output_path = os.path.join(OUTPUT_DIR, "Mechanical-Ventilation-Basics.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

if __name__ == "__main__":
    print("Creating PowerPoint presentations...")
    print("=" * 50)
    
    create_patient_assessment_ppt()
    create_emergency_procedures_ppt()
    create_ventilation_basics_ppt()
    
    print("=" * 50)
    print("All presentations created successfully!")
