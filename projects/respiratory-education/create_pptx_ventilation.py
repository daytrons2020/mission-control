#!/usr/bin/env python3
"""
Create comprehensive PowerPoint presentations for Mechanical Ventilation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/root/.openclaw/workspace/projects/respiratory-education/output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    return slide

def add_content_slide(prs, title, content, notes=""):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    lines = content.strip().split('\n')
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        is_header = line.startswith('**') and line.endswith('**')
        is_subheader = line.startswith('##') or (line.startswith('**') and not is_header)
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        clean_line = line.replace('**', '').replace('##', '').replace('#', '').strip()
        if clean_line.startswith('- ') or clean_line.startswith('* '):
            clean_line = clean_line[2:]
            p.level = 1
        elif clean_line.startswith('1.') or clean_line.startswith('2.') or clean_line.startswith('3.'):
            p.level = 1
        
        p.text = clean_line
        if is_header or is_subheader:
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 102, 153)
            p.level = 0
        else:
            p.font.size = Pt(16)
            p.level = 1 if p.level == 1 else 0
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def create_ventilation_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, "Mechanical Ventilation Management", 
                   "Comprehensive Guide to Ventilator Management\n\nPresenter: [Name], RRT\nDate: [Date]\nDepartment: Respiratory Care")
    
    # Learning Objectives
    add_content_slide(prs, "Learning Objectives", 
                     """By the end of this session, participants will be able to:
- Describe the components of a mechanical ventilator and breathing circuit
- Select appropriate ventilator modes based on patient condition
- Apply lung-protective ventilation strategies
- Interpret ventilator graphics and waveforms
- Troubleshoot common ventilator alarms and patient-ventilator asynchrony
- Develop and implement ventilator weaning protocols""",
                     "Mechanical ventilation is a core skill for respiratory therapists. This presentation covers essential concepts.")
    
    # Indications
    add_content_slide(prs, "Indications for Mechanical Ventilation",
                     """**Primary Indications:**

**Apnea or Respiratory Arrest**
- No spontaneous breathing
- Immediate intervention required

**Acute Respiratory Failure**
- Hypoxemic: PaO2 <60 mmHg on high FIO2
- Hypercapnic: pH <7.30 with elevated PaCO2
- Increased work of breathing

**Airway Protection**
- Altered mental status (GCS ≤8)
- Inability to protect airway
- Risk of aspiration

**Prophylactic/Perioperative**
- Major surgery
- Hemodynamic instability
- Severe metabolic acidosis""",
                     "Indications for mechanical ventilation include respiratory failure, airway protection, and perioperative support.")
    
    # Ventilator Components
    add_content_slide(prs, "Ventilator Components",
                     """**The Mechanical Ventilator System**

**Control System:**
- Microprocessor control
- User interface
- Monitoring displays
- Alarm systems

**Pneumatic System:**
- Gas source (wall oxygen, medical air)
- Pressure regulators
- Flow controllers
- Valves

**Breathing Circuit:**
- Inspiratory limb
- Expiratory limb
- Y-connector
- Humidifier
- HME or heated humidifier

**Patient Interface:**
- Endotracheal tube
- Tracheostomy tube
- Mask (for NIV)""",
                     "Understanding ventilator components is essential for troubleshooting and proper setup.")
    
    # Ventilator Modes Overview
    add_content_slide(prs, "Ventilator Modes: Overview",
                     """**Classification by Control Variable**

**Volume Control (VC)**
- Delivers set tidal volume
- Pressure varies with compliance/resistance
- Guaranteed minute ventilation
- Examples: VC-AC, VC-SIMV

**Pressure Control (PC)**
- Delivers set inspiratory pressure
- Volume varies with compliance/resistance
- Decelerating flow pattern
- Examples: PC-AC, PC-SIMV, PRVC

**Pressure Support (PS)**
- Patient-triggered, pressure-limited
- Flow-cycled
- Used for spontaneous breathing""",
                     "Ventilator modes are classified by what variable is controlled: volume, pressure, or support.")
    
    # Assist/Control
    add_content_slide(prs, "Ventilator Modes: Assist/Control (A/C)",
                     """**Volume Assist/Control**

**Characteristics:**
- All breaths are volume-controlled
- Patient can trigger breaths
- If no trigger, ventilator delivers at set rate
- Every breath has same volume

**Settings:**
- Tidal volume (VT)
- Respiratory rate (RR)
- FIO2
- PEEP
- Flow rate
- Inspiratory time or I:E ratio

**Advantages:**
- Guaranteed minute ventilation
- Good for patients with no spontaneous effort

**Disadvantages:**
- Can cause asynchrony if patient wants different pattern
- Risk of barotrauma if patient fights ventilator""",
                     "A/C mode ensures every breath delivers the set volume. It's commonly used for patients requiring full support.")
    
    # SIMV
    add_content_slide(prs, "Ventilator Modes: SIMV",
                     """**Synchronized Intermittent Mandatory Ventilation**

**Characteristics:**
- Set number of mandatory breaths (volume or pressure)
- Patient can breathe spontaneously between mandatory breaths
- Spontaneous breaths can be unsupported or pressure-supported

**Settings:**
- Mandatory rate
- Tidal volume (if VC) or pressure (if PC)
- Pressure support for spontaneous breaths
- FIO2, PEEP

**Advantages:**
- Allows spontaneous breathing
- May reduce sedation needs
- Theoretically prevents respiratory muscle atrophy

**Disadvantages:**
- Can increase work of breathing
- May prolong weaning
- Not proven superior to A/C""",
                     "SIMV allows spontaneous breathing between mandatory breaths. Its benefits over A/C are not well established.")
    
    # Pressure Support
    add_content_slide(prs, "Ventilator Modes: Pressure Support",
                     """**Spontaneous Mode with Support**

**Characteristics:**
- Patient triggers every breath
- Set pressure support level
- Flow-cycled (cycles off when flow drops to % of peak)
- No set rate or tidal volume

**Settings:**
- Pressure support level (cmH2O)
- PEEP
- FIO2
- Trigger sensitivity
- Flow cycle criteria

**Advantages:**
- Patient controls rate and tidal volume
- Comfortable for awake patients
- Good for weaning

**Disadvantages:**
- No guaranteed minute ventilation
- Apnea alarm essential
- Not for patients with weak drive""",
                     "Pressure support is comfortable for patients and commonly used for weaning from mechanical ventilation.")
    
    # Advanced Modes
    add_content_slide(prs, "Ventilator Modes: Advanced Modes",
                     """**Pressure Regulated Volume Control (PRVC/VC+)**

**Characteristics:**
- Pressure-controlled breaths
- Volume-targeted
- Automatically adjusts pressure to deliver set volume

**Advantages:**
- Decelerating flow (better distribution)
- Volume guarantee
- Automatic compliance compensation

**Airway Pressure Release Ventilation (APRV)**

**Characteristics:**
- Two pressure levels (Phigh and Plow)
- Spontaneous breathing allowed at both levels
- Time-cycled release

**Advantages:**
- Good oxygenation
- Allows spontaneous breathing
- Lower peak pressures""",
                     "Advanced modes combine benefits of pressure and volume control. APRV is useful for refractory hypoxemia.")
    
    # Lung-Protective Ventilation
    add_content_slide(prs, "Lung-Protective Ventilation",
                     """**ARDSNet Protocol Principles**

**Tidal Volume:**
- 6 ml/kg ideal body weight
- Range: 4-8 ml/kg
- Calculate IBW:
  - Men: 50 + 2.3 × (height in inches - 60)
  - Women: 45.5 + 2.3 × (height in inches - 60)

**Plateau Pressure:**
- Goal: <30 cmH2O
- If >30: decrease VT to minimum 4 ml/kg

**PEEP:**
- Start at 5 cmH2O
- Titrate using PEEP/FIO2 tables
- Balance oxygenation vs. hemodynamics

**Permissive Hypercapnia:**
- Accept PaCO2 up to 60 if pH ≥7.25
- Don't increase VT to normalize CO2""",
                     "Lung-protective ventilation reduces ventilator-induced lung injury. The ARDSNet protocol is the standard of care.")
    
    # Initial Setup
    add_content_slide(prs, "Initial Ventilator Setup",
                     """**Standard Starting Points**

**Volume A/C Mode:**
| Parameter | Setting | Notes |
|-----------|---------|-------|
| VT | 6-8 ml/kg IBW | Use IBW, not actual weight |
| Rate | 12-16/min | Adjust for pH/CO2 |
| FIO2 | 100% | Titrate down quickly |
| PEEP | 5 cmH2O | Minimum to prevent atelectasis |
| Flow | 60 L/min | Adequate for most adults |
| I:E ratio | 1:2 to 1:3 | Avoid inverse ratio |

**Post-Intubation:**
- Check bilateral breath sounds
- Verify ET tube position (depth marking)
- Obtain chest X-ray
- Check ABG in 30 minutes""",
                     "Standard initial settings provide a starting point. Always adjust based on patient response and ABG results.")
    
    # Ventilator Graphics
    add_content_slide(prs, "Ventilator Graphics",
                     """**Understanding Waveforms**

**Pressure-Time Waveform:**
- Shows pressure changes during breath
- Peak pressure: Total resistance + compliance
- Plateau pressure: Static compliance only
- PEEP: Baseline pressure

**Flow-Time Waveform:**
- Square wave (constant flow) - volume control
- Decelerating wave - pressure control
- Expiratory flow shows airway obstruction

**Volume-Time Waveform:**
- Shows volume delivery and exhalation
- Should return to zero
- If not, air trapping present""",
                     "Ventilator graphics provide valuable information about patient-ventilator interaction and lung mechanics.")
    
    # Graphics Interpretation
    add_content_slide(prs, "Ventilator Graphics Interpretation",
                     """**Common Patterns**

**Normal:**
- Smooth, consistent waveforms
- Volume returns to baseline
- Stable pressures

**Airway Obstruction:**
- Prolonged expiratory flow
- Scalloped expiratory flow waveform
- Auto-PEEP

**Asynchrony:**
- Pressure waveform notches (triggering issues)
- Flow waveform abnormalities
- Double triggering

**Bronchodilator Response:**
- Decreased peak pressure
- Improved expiratory flow
- Reduced auto-PEEP""",
                     "Recognizing waveform patterns helps identify patient-ventilator asynchrony and response to treatment.")
    
    # Patient-Ventilator Asynchrony
    add_content_slide(prs, "Patient-Ventilator Asynchrony",
                     """**Types and Solutions**

**Trigger Asynchrony:**
- **Problem:** Patient effort doesn't trigger breath
- **Cause:** Insensitive trigger, auto-PEEP
- **Solution:** Increase sensitivity, reduce auto-PEEP

**Flow Asynchrony:**
- **Problem:** Set flow doesn't match patient demand
- **Cause:** Flow too low for patient's needs
- **Solution:** Increase flow or switch to pressure control

**Cycle Asynchrony:**
- **Problem:** Inspiration ends too early or late
- **Cause:** Mismatch between neural and mechanical inspiration
- **Solution:** Adjust flow cycle % or inspiratory time

**Double Triggering:**
- **Problem:** Two breaths in rapid succession
- **Cause:** Insufficient tidal volume, short inspiratory time
- **Solution:** Increase VT or adjust flow""",
                     "Asynchrony increases work of breathing and patient discomfort. Identifying the type guides correction.")
    
    # High Pressure Alarm
    add_content_slide(prs, "Ventilator Alarms: High Pressure",
                     """**Causes:**
- Secretions in airway
- Kinked ETT or circuit
- Bronchospasm
- Decreased lung compliance
- Patient biting ETT
- Water in circuit
- Coughing or fighting ventilator

**Response:**
1. Check patient first - look, listen, feel
2. Suction if indicated
3. Check circuit for kinks or water
4. Auscultate for bronchospasm
5. Check ETT position and patency
6. Consider sedation if patient-ventilator dyssynchrony
7. If unresolved, disconnect and bag manually""",
                     "Always check the patient first when alarms sound. Systematic troubleshooting resolves most issues.")
    
    # Low Pressure Alarm
    add_content_slide(prs, "Ventilator Alarms: Low Pressure/Volume",
                     """**Causes:**
- Circuit disconnect
- Cuff leak
- Patient self-extubation
- Leak in circuit

**Response:**
1. Check patient immediately
2. Verify ETT position (depth marking)
3. Check cuff pressure (20-30 cmH2O)
4. Check all connections
5. Listen for leak around ETT
6. If extubated, bag-valve-mask and call for help

**Low PEEP/CPAP Alarm:**
- Large leak in system
- Circuit disconnect
- Inadequate flow""",
                     "Low pressure alarms may indicate a critical disconnection. Respond immediately to check the patient.")
    
    # Auto-PEEP
    add_content_slide(prs, "Auto-PEEP (Intrinsic PEEP)",
                     """**Definition:**
Positive pressure remaining in alveoli at end-expiration due to incomplete emptying

**Causes:**
- Airway obstruction (asthma, COPD)
- High minute ventilation
- Short expiratory time
- High respiratory rate

**Detection:**
- Expiratory flow doesn't return to zero
- End-expiratory port occlusion shows pressure
- Need to trigger against positive pressure

**Management:**
- Increase expiratory time (decrease rate or inspiratory time)
- Bronchodilators
- Reduce minute ventilation (if appropriate)
- Apply extrinsic PEEP (80% of auto-PEEP)""",
                     "Auto-PEEP increases work of breathing and can cause hemodynamic compromise. It's common in obstructive diseases.")
    
    # Weaning
    add_content_slide(prs, "Weaning from Mechanical Ventilation",
                     """**Daily Spontaneous Awakening and Breathing Trials**

**Readiness Criteria:**
- Resolution of acute illness
- Hemodynamic stability (minimal or no vasopressors)
- Oxygenation: PaO2/FIO2 >150-200, PEEP ≤5-8
- Cough and airway reflexes present
- No pending procedures requiring sedation

**Spontaneous Breathing Trial (SBT):**
- CPAP 5 cmH2O, or
- T-piece, or
- Low pressure support (5-8 cmH2O)
- Duration: 30-120 minutes

**SBT Success Criteria:**
- RR <35
- SpO2 >90%
- HR <140 or <20% increase
- No hemodynamic instability
- No respiratory distress""",
                     "Daily SBTs identify patients ready for extubation. Most patients who fail do so within the first few minutes.")
    
    # Weaning Parameters
    add_content_slide(prs, "Weaning Parameters",
                     """**Objective Measures**

**Rapid Shallow Breathing Index (RSBI):**
- RR / VT (in liters)
- <105: Likely to succeed
- >105: Likely to fail

**Negative Inspiratory Force (NIF):**
- <-20 to -30 cmH2O
- Measures inspiratory muscle strength

**Other Parameters:**
- Spontaneous VT >5 ml/kg
- Vital capacity >10-15 ml/kg
- Minute ventilation <10 L/min
- Maximum voluntary ventilation >2× resting VE

**Note:** These predict success better than failure. Clinical assessment is most important.""",
                     "Weaning parameters help predict success but should not replace clinical judgment.")
    
    # Extubation
    add_content_slide(prs, "Extubation",
                     """**Procedure and Considerations**

**Pre-Extubation:**
- Suction ETT and oropharynx
- Suction below vocal cords if possible
- Position patient upright
- Explain procedure to patient
- Have reintubation equipment ready

**Procedure:**
1. Deflate cuff
2. Give positive pressure breath
3. Remove tube at peak inspiration
4. Apply oxygen (mask or nasal cannula)
5. Encourage cough

**Post-Extubation:**
- Monitor closely for stridor
- Humidified oxygen
- Consider noninvasive ventilation if high-risk
- Monitor for respiratory fatigue""",
                     "Proper preparation and technique reduce extubation failure. Monitor closely for signs of upper airway obstruction.")
    
    # Post-Extubation Management
    add_content_slide(prs, "Post-Extubation Management",
                     """**Preventing Reintubation**

**Risk Factors for Extubation Failure:**
- Prolonged intubation
- Weak cough
- Excess secretions
- Upper airway edema
- CHF
- Neuromuscular disease

**Preventive Strategies:**
- Noninvasive ventilation (prophylactic in high-risk)
- Aggressive pulmonary toilet
- Humidification
- Early mobilization
- Treat underlying causes

**Reintubation Criteria:**
- Respiratory failure
- Inability to protect airway
- Severe stridor not responding to treatment
- Hemodynamic instability""",
                     "Prophylactic NIV may prevent reintubation in high-risk patients. Early intervention is key.")
    
    # Case Study 1
    add_content_slide(prs, "Case Study 1",
                     """**Patient:** Mr. Smith, 65-year-old male
- Post-operative day 2 after abdominal surgery
- Intubated for respiratory failure
- Current settings: VC-AC, VT 500, Rate 14, FIO2 40%, PEEP 5

**Current Status:**
- Peak pressure: 35 cmH2O
- Plateau pressure: 28 cmH2O
- SpO2: 94%
- Patient appears uncomfortable, asynchronous with ventilator

**Questions:**
1. Are these pressures acceptable?
2. What might be causing the asynchrony?
3. What adjustments would you make?""",
                     "This patient shows signs of patient-ventilator asynchrony post-surgery.")
    
    # Case Study 1 Discussion
    add_content_slide(prs, "Case Study 1 Discussion",
                     """**Pressure Assessment:**
- Plateau 28: Acceptable (<30)
- Peak-plateau difference: 7 cmH2O (normal airway resistance)
- Pressures are acceptable

**Asynchrony Causes:**
- Pain/anxiety post-surgery
- Inadequate sedation
- Flow asynchrony (set flow may not match demand)
- Possible inadequate tidal volume for patient

**Interventions:**
1. Assess pain and sedation needs
2. Consider pressure support or PC mode for comfort
3. Check for auto-PEEP
4. If flow asynchrony suspected, increase flow or switch to PC
5. Consider spontaneous breathing trial if criteria met""",
                     "Addressing pain, sedation, and ventilator settings can improve comfort and synchrony.")
    
    # Case Study 2
    add_content_slide(prs, "Case Study 2",
                     """**Patient:** Mrs. Johnson, 58-year-old female
- ARDS secondary to sepsis
- Day 5 of mechanical ventilation
- Current settings: VC-AC, VT 400 (6 ml/kg), Rate 22, FIO2 70%, PEEP 12

**Current Status:**
- Plateau pressure: 32 cmH2O
- SpO2: 88%
- Chest X-ray: Bilateral infiltrates

**Questions:**
1. What is the immediate concern?
2. How would you adjust the ventilator?
3. What other interventions should be considered?""",
                     "This patient with ARDS has high plateau pressure and hypoxemia despite high support.")
    
    # Case Study 2 Discussion
    add_content_slide(prs, "Case Study 2 Discussion",
                     """**Immediate Concern:**
- Plateau pressure 32 exceeds safe limit (<30)
- Risk of ventilator-induced lung injury
- Severe hypoxemia despite high support

**Ventilator Adjustments:**
1. Decrease VT to 300-350 ml (4-5 ml/kg)
2. Check and document driving pressure (Pplat - PEEP)
3. Increase PEEP per ARDSNet table
4. Titrate FIO2 to maintain SpO2 88-95%

**Other Interventions:**
- Prone positioning if PaO2/FIO2 <150
- Consider neuromuscular blockade if severe
- Evaluate for ECMO if refractory
- Treat underlying sepsis
- Conservative fluid management""",
                     "In ARDS, lung-protective ventilation is essential. Consider adjunctive therapies for refractory hypoxemia.")
    
    # Key Takeaways
    add_content_slide(prs, "Key Takeaways",
                     """1. **Select Mode Based on Patient Needs**
   - Volume control for guaranteed ventilation
   - Pressure control for comfort
   - Pressure support for weaning

2. **Always Use Lung-Protective Strategy**
   - 6 ml/kg IBW, Pplat <30
   - Accept permissive hypercapnia if needed

3. **Graphics Tell the Story**
   - Learn to recognize normal and abnormal patterns
   - Waveforms guide troubleshooting

4. **Treat the Patient, Not the Numbers**
   - Clinical assessment trumps parameters
   - Comfort and synchrony matter

5. **Wean Early and Daily**
   - Daily SBTs for eligible patients
   - Use spontaneous awakening trials""",
                     "Mechanical ventilation requires balancing multiple factors. Always prioritize lung protection and patient comfort.")
    
    # References
    add_content_slide(prs, "References",
                     """1. ARDSNet. Ventilation with lower tidal volumes as compared with traditional tidal volumes for acute lung injury and the acute respiratory distress syndrome. N Engl J Med. 2000;342(18):1301-1308.

2. AARC Clinical Practice Guidelines: Mechanical Ventilation of Mechanically Ventilated Patients. Respir Care. 2022.

3. Girard TD, et al. Efficacy and safety of a paired sedation and ventilator weaning protocol for mechanically ventilated patients in intensive care. Lancet. 2008;371(9607):126-134.

4. Papazian L, et al. Formal guidelines: management of acute respiratory distress syndrome. Ann Intensive Care. 2019;9(1):69.

5. Tobin MJ. Principles and Practice of Mechanical Ventilation. 4th ed. McGraw-Hill; 2018.""",
                     "The ARDSNet trial established lung-protective ventilation as the standard of care.")
    
    add_title_slide(prs, "Thank You", "Respiratory Care Department\n\nExcellence in Mechanical Ventilation Support")
    
    output_path = os.path.join(OUTPUT_DIR, "03-Mechanical-Ventilation.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_ventilation_presentation()
