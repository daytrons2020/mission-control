#!/usr/bin/env python3
"""
Respiratory Education PowerPoint and Word Document Generator
Creates professional slide decks and detailed Word documents from markdown content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
import re
import os

# Brand colors
PRIMARY_COLOR = RGBColor(0x1E, 0x5A, 0xA8)  # Professional blue
ACCENT_COLOR = RGBColor(0x2E, 0x8B, 0x57)   # Forest green
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)     # Dark gray
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)       # Light gray background

class RespiratoryEducationGenerator:
    def __init__(self, output_dir):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
    
    def parse_markdown(self, content):
        """Parse markdown content into slides/sections."""
        slides = []
        current_slide = {"title": "", "content": [], "speaker_notes": ""}
        
        lines = content.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i]
            
            # Title slide detection
            if line.startswith('# ') and i < len(lines) - 1 and '---' in lines[i-1:i+1]:
                if current_slide["title"] or current_slide["content"]:
                    slides.append(current_slide)
                current_slide = {"title": line[2:].strip(), "content": [], "speaker_notes": "", "type": "title"}
            
            # Section header (new slide)
            elif line.startswith('# ') or line.startswith('## '):
                if current_slide["title"] or current_slide["content"]:
                    slides.append(current_slide)
                title = line[line.find(' ')+1:].strip()
                current_slide = {"title": title, "content": [], "speaker_notes": "", "type": "content"}
            
            # Speaker notes
            elif '**Speaker Notes:**' in line:
                notes = []
                i += 1
                while i < len(lines) and not lines[i].startswith('#') and not lines[i].startswith('---'):
                    if lines[i].strip():
                        notes.append(lines[i].strip())
                    i += 1
                current_slide["speaker_notes"] = '\n'.join(notes)
                continue
            
            # Content lines
            elif line.strip() and not line.startswith('---') and not line.startswith('<!--'):
                # Clean up markdown formatting
                clean_line = line.strip()
                clean_line = re.sub(r'\*\*', '', clean_line)  # Remove bold markers
                clean_line = re.sub(r'\*', '', clean_line)    # Remove italic markers
                clean_line = re.sub(r'`', '', clean_line)     # Remove code markers
                current_slide["content"].append(clean_line)
            
            i += 1
        
        if current_slide["title"] or current_slide["content"]:
            slides.append(current_slide)
        
        return slides
    
    def create_presentation(self, title, slides, filename):
        """Create a PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for i, slide_data in enumerate(slides):
            if i == 0 or slide_data.get("type") == "title":
                slide_layout = prs.slide_layouts[6]  # Blank
            else:
                slide_layout = prs.slide_layouts[6]  # Blank
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Add background shape
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            bg.line.fill.background()
            
            # Title bar for content slides
            if i > 0:
                title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
                title_bar.fill.solid()
                title_bar.fill.fore_color.rgb = PRIMARY_COLOR
                title_bar.line.fill.background()
            
            # Add title
            if slide_data["title"]:
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_data["title"]
                p.font.size = Pt(32 if i == 0 else 28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i > 0 else PRIMARY_COLOR
                p.alignment = PP_ALIGN.LEFT
            
            # Add content
            if slide_data["content"]:
                content_text = '\n'.join(slide_data["content"])
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
                tf = content_box.text_frame
                tf.word_wrap = True
                
                # Split content into paragraphs
                paragraphs = content_text.split('\n')
                for j, para_text in enumerate(paragraphs[:15]):  # Limit to avoid overflow
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = para_text[:200]  # Limit line length
                    p.font.size = Pt(18)
                    p.font.color.rgb = TEXT_COLOR
                    p.space_after = Pt(8)
                    
                    # Indent bullet points
                    if para_text.strip().startswith('-') or para_text.strip().startswith('•'):
                        p.level = 1
                        p.font.size = Pt(16)
            
            # Add footer with slide number
            footer = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
            tf = footer.text_frame
            p = tf.paragraphs[0]
            p.text = f"{i + 1}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
            p.alignment = PP_ALIGN.RIGHT
            
            # Add speaker notes
            if slide_data.get("speaker_notes"):
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data["speaker_notes"][:1000]  # Limit length
        
        # Save presentation
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)
        print(f"Created PowerPoint: {output_path}")
        return output_path
    
    def create_word_document(self, title, content, filename):
        """Create a Word document from outline content."""
        doc = Document()
        
        # Set up styles
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Calibri'
        font.size = DocxPt(11)
        
        # Add title
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add content
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if not line or line.startswith('---') or line.startswith('<!--'):
                continue
            
            # Clean markdown
            line = re.sub(r'\*\*', '', line)
            line = re.sub(r'\*', '', line)
            line = re.sub(r'`', '', line)
            
            # Headers
            if line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.startswith('#### '):
                doc.add_heading(line[5:], level=4)
            # Speaker notes section
            elif '**Speaker Notes:**' in line:
                doc.add_heading('Speaker Notes', level=3)
            # Table rows
            elif '|' in line and not line.startswith('|--'):
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                if cells:
                    p = doc.add_paragraph()
                    p.add_run(' | '.join(cells)).italic = True
            # Regular content
            elif line:
                if line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line[2:], style='List Bullet')
                elif line[0].isdigit() and '. ' in line[:4]:
                    doc.add_paragraph(line[line.find('. ')+2:], style='List Number')
                else:
                    doc.add_paragraph(line)
        
        # Save document
        output_path = os.path.join(self.output_dir, filename)
        doc.save(output_path)
        print(f"Created Word document: {output_path}")
        return output_path
    
    def process_file(self, md_path, ppt_filename, doc_filename=None):
        """Process a markdown file to create both PowerPoint and Word document."""
        with open(md_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Extract title
        title_match = re.search(r'^# (.+)$', content, re.MULTILINE)
        title = title_match.group(1) if title_match else "Respiratory Education"
        
        # Create PowerPoint
        slides = self.parse_markdown(content)
        ppt_path = self.create_presentation(title, slides, ppt_filename)
        
        # Create Word document if requested
        doc_path = None
        if doc_filename:
            doc_path = self.create_word_document(title, content, doc_filename)
        
        return ppt_path, doc_path


def main():
    output_dir = "/root/.openclaw/workspace/projects/respiratory-education/output"
    templates_dir = "/root/.openclaw/workspace/projects/respiratory-education/templates"
    
    generator = RespiratoryEducationGenerator(output_dir)
    
    # Process files
    files_to_process = [
        # (markdown_file, ppt_filename, doc_filename)
        ("01-anatomy-physiology.md", "01-Anatomy-Physiology.pptx", "01-Anatomy-Physiology.docx"),
        ("02-copd-management.md", "02-COPD-Management.pptx", "02-COPD-Management.docx"),
        ("03-mechanical-ventilation.md", "03-Mechanical-Ventilation.pptx", "03-Mechanical-Ventilation.docx"),
        ("04-patient-assessment.md", "04-Patient-Assessment.pptx", "04-Patient-Assessment.docx"),
        ("05-emergency-procedures.md", "05-Emergency-Procedures.pptx", "05-Emergency-Procedures.docx"),
        ("mechanical-ventilation-basics.md", "Mechanical-Ventilation-Basics.pptx", None),
    ]
    
    for md_file, ppt_file, doc_file in files_to_process:
        md_path = os.path.join(templates_dir, md_file)
        if os.path.exists(md_path):
            print(f"\nProcessing {md_file}...")
            generator.process_file(md_path, ppt_file, doc_file)
        else:
            print(f"Warning: {md_path} not found")
    
    print("\n✓ All files processed successfully!")


if __name__ == "__main__":
    main()
