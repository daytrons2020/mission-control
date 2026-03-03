#!/usr/bin/env python3
"""
Create comprehensive PowerPoint presentations for COPD Management
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/root/.openclaw/workspace/projects/respiratory-education/output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    return slide

def add_content_slide(prs, title, content, notes=""):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    lines = content.strip().split('\n')
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        is_header = line.startswith('**') and line.endswith('**')
        is_subheader = line.startswith('##') or (line.startswith('**') and not is_header)
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        clean_line = line.replace('**', '').replace('##', '').replace('#', '').strip()
        if clean_line.startswith('- ') or clean_line.startswith('* '):
            clean_line = clean_line[2:]
            p.level = 1
        elif clean_line.startswith('1.') or clean_line.startswith('2.') or clean_line.startswith('3.'):
            p.level = 1
        
        p.text = clean_line
        if is_header or is_subheader:
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 102, 153)
            p.level = 0
        else:
            p.font.size = Pt(16)
            p.level = 1 if p.level == 1 else 0
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def create_copd_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, "Chronic Obstructive Pulmonary Disease (COPD)", 
                   "Pathophysiology, Management, and Patient Education\n\nPresenter: [Name], RRT\nDate: [Date]\nDepartment: Respiratory Care")
    
    # Learning Objectives
    add_content_slide(prs, "Learning Objectives", 
                     """By the end of this session, participants will be able to:
- Define COPD and differentiate between emphysema and chronic bronchitis
- Explain the pathophysiology of airflow limitation in COPD
- Identify the stages of COPD using GOLD criteria
- Describe pharmacological and non-pharmacological management strategies
- Manage acute exacerbations including oxygen therapy and NIV
- Develop patient education strategies for smoking cessation and self-management""",
                     "COPD is a major cause of morbidity and mortality worldwide. Understanding its management is essential for respiratory therapists.")
    
    # Definition
    add_content_slide(prs, "Definition of COPD",
                     """**Global Initiative for Chronic Obstructive Lung Disease (GOLD)**

**Definition:**
"COPD is a common, preventable, and treatable disease characterized by persistent respiratory symptoms and airflow limitation due to airway and/or alveolar abnormalities usually caused by significant exposure to noxious particles or gases."

**Key Characteristics:**
- **Preventable:** Primary prevention through smoking cessation
- **Treatable:** Not curable, but manageable
- **Persistent:** Chronic, progressive condition
- **Airflow limitation:** Not fully reversible""",
                     "The GOLD definition emphasizes that COPD is preventable and treatable, even though not curable. Smoking cessation is the most important intervention.")
    
    # Types of COPD
    add_content_slide(prs, "Types of COPD: Emphysema vs. Chronic Bronchitis",
                     """**Emphysema ("Pink Puffer")**
- Destruction of alveolar walls
- Loss of elastic recoil
- Increased compliance
- Panacinar or centriacinar patterns
- Dyspnea predominant symptom
- Thin, non-cyanotic appearance

**Chronic Bronchitis ("Blue Bloater")**
- Productive cough for 3+ months/year for 2+ years
- Mucus hypersecretion
- Airway inflammation
- Cyanosis and edema common
- Recurrent infections
- Obesity common

**Note:** Most patients have features of both.""",
                     "The classic phenotypes are less distinct in practice. Most patients exhibit features of both emphysema and chronic bronchitis.")
    
    # Epidemiology
    add_content_slide(prs, "Epidemiology",
                     """**Global Statistics:**
- Third leading cause of death worldwide
- 328 million people affected globally
- Prevalence increasing in developing nations

**United States:**
- 16 million diagnosed (estimated 12 million undiagnosed)
- $50 billion annual healthcare costs
- Leading cause of hospital readmissions

**Risk Factors:**
- Tobacco smoking (85-90% of cases)
- Alpha-1 antitrypsin deficiency
- Occupational exposures
- Biomass fuel exposure
- Air pollution
- History of childhood respiratory infections""",
                     "COPD is a major public health burden. Many cases remain undiagnosed until significant lung function is lost.")
    
    # Pathophysiology - Airways
    add_content_slide(prs, "Pathophysiology: Airway Changes",
                     """**Chronic Inflammation and Remodeling**

**Small Airways (<2mm)**
- Inflammation of bronchioles
- Fibrosis and narrowing
- Excess mucus production
- Loss of tethering support

**Large Airways**
- Mucus gland hypertrophy
- Goblet cell hyperplasia
- Squamous metaplasia
- Cilia dysfunction

**Result:**
- Fixed airflow limitation
- Increased airway resistance
- Premature airway closure
- Air trapping""",
                     "Chronic inflammation leads to structural changes in both small and large airways, resulting in fixed airflow limitation.")
    
    # Pathophysiology - Parenchyma
    add_content_slide(prs, "Pathophysiology: Parenchymal Changes",
                     """**Emphysema Patterns**

**Centriacinar (Centrilobular)**
- Central portion of acinus affected
- Upper lobes predominate
- Associated with smoking
- Most common type (95%)

**Panacinar (Panlobular)**
- Entire acinus affected
- Lower lobes predominate
- Associated with alpha-1 antitrypsin deficiency

**Pathological Process:**
- Imbalance of proteases and antiproteases
- Elastin destruction
- Loss of alveolar attachments
- Airspace enlargement""",
                     "Emphysema patterns differ based on etiology. Understanding these patterns helps identify underlying causes.")
    
    # Pathophysiology - Physiological
    add_content_slide(prs, "Pathophysiology: Physiological Consequences",
                     """**Ventilation-Perfusion (V/Q) Mismatch**

**Mechanisms:**
- Destruction of alveolar-capillary membrane
- Hypoxic vasoconstriction
- Reduced surface area for gas exchange

**Consequences:**
- Hypoxemia (low PaO2)
- Hypercapnia (elevated PaCO2) in advanced disease
- Pulmonary hypertension
- Cor pulmonale (right heart failure)

**Compensatory Mechanisms:**
- Increased respiratory rate
- Accessory muscle use
- Barrel chest (increased AP diameter)""",
                     "V/Q mismatch is the primary cause of hypoxemia in COPD. As disease progresses, hypercapnia and pulmonary hypertension develop.")
    
    # Clinical Presentation
    add_content_slide(prs, "Clinical Presentation",
                     """**Cardinal Symptoms:**
- Progressive dyspnea
- Chronic cough
- Sputum production
- Wheezing

**Physical Examination:**
- Prolonged expiratory phase
- Decreased breath sounds
- Wheezing or rhonchi
- Use of accessory muscles
- Pursed-lip breathing
- Barrel chest
- Digital clubbing (late sign)

**Advanced Disease:**
- Weight loss
- Cachexia
- Peripheral edema
- Cyanosis""",
                     "Dyspnea that progressively worsens is the hallmark symptom. Physical findings correlate with disease severity.")
    
    # Diagnostic Testing
    add_content_slide(prs, "Diagnostic Testing",
                     """**Pulmonary Function Tests**

**Spirometry (Required for Diagnosis)**
- Post-bronchodilator FEV1/FVC < 0.70 confirms airflow limitation
- FEV1 % predicted determines severity

**Lung Volumes**
- Increased total lung capacity (TLC)
- Increased residual volume (RV)
- Increased RV/TLC ratio (air trapping)

**Diffusion Capacity (DLCO)**
- Reduced in emphysema
- Normal or slightly reduced in chronic bronchitis

**Arterial Blood Gas**
- Hypoxemia (common)
- Hypercapnia (in advanced disease)
- Respiratory acidosis (chronic or acute-on-chronic)""",
                     "Spirometry is essential for diagnosis. The FEV1/FVC ratio defines airflow limitation, and FEV1 determines severity.")
    
    # GOLD Classification
    add_content_slide(prs, "GOLD Classification",
                     """**Severity Grading (Post-Bronchodilator FEV1)**

| Stage | FEV1 % Predicted | Description |
|-------|-----------------|-------------|
| **GOLD 1** | ≥80% | Mild |
| **GOLD 2** | 50-79% | Moderate |
| **GOLD 3** | 30-49% | Severe |
| **GOLD 4** | <30% | Very Severe |

**Note:** Symptoms and exacerbation history also guide treatment (ABE assessment).""",
                     "GOLD classification guides treatment decisions. However, symptoms and exacerbation history are equally important for management.")
    
    # GOLD Assessment Framework
    add_content_slide(prs, "GOLD Assessment Framework",
                     """**Combined Assessment (2024 Update)**

**Symptom Assessment:**
- mMRC Dyspnea Scale (0-4)
- CAT Score (0-40)
- Group A: 0-1 symptoms
- Group B: ≥2 symptoms or mMRC ≥2
- Group E: Exacerbation history

**Exacerbation History:**
- Group E: ≥2 moderate exacerbations or ≥1 hospitalization

**Combined Groups:**
- A, B, E (with or without exacerbation risk)""",
                     "The ABE assessment framework combines symptoms and exacerbation risk to guide initial pharmacological treatment.")
    
    # Pharmacological Management - Bronchodilators
    add_content_slide(prs, "Pharmacological Management: Bronchodilators",
                     """**Mainstay of Therapy**

**Short-Acting Bronchodilators (Rescue)**
- SABA: Albuterol, levalbuterol
- SAMA: Ipratropium bromide
- Use: PRN for symptom relief

**Long-Acting Bronchodilators (Maintenance)**
- LABA: Salmeterol, formoterol, indacaterol
- LAMA: Tiotropium, umeclidinium, glycopyrrolate
- Use: Daily for symptom control

**Combination Therapy**
- LABA/LAMA preferred over single agents
- SABA/SAMA for acute symptoms""",
                     "Bronchodilators are the foundation of COPD treatment. Long-acting agents are used for maintenance, short-acting for rescue.")
    
    # Pharmacological Management - ICS
    add_content_slide(prs, "Pharmacological Management: Inhaled Corticosteroids",
                     """**Indications:**
- History of exacerbations
- Asthma-COPD overlap
- Eosinophil count ≥300/μL

**Risks:**
- Pneumonia (increased risk)
- Oral candidiasis
- Hoarseness

**Triple Therapy:**
- ICS/LABA/LAMA for frequent exacerbators
- Examples: Trelegy, Breztri

**Oral Corticosteroids:**
- Acute exacerbations only
- Not for chronic use""",
                     "ICS should be reserved for patients with frequent exacerbations or high eosinophil counts due to pneumonia risk.")
    
    # Other Medications
    add_content_slide(prs, "Other Medications",
                     """**Phosphodiesterase-4 Inhibitors**
- Roflumilast (Daliresp)
- For severe COPD with chronic bronchitis
- Reduces exacerbations

**Mucolytics**
- Acetylcysteine, carbocysteine
- May reduce exacerbations
- Limited evidence

**Antibiotics (Prophylactic)**
- Azithromycin (250mg daily)
- For frequent exacerbators
- Monitor QT interval and hearing

**Alpha-1 Antitrypsin Augmentation**
- For AAT deficiency
- Weekly IV infusions""",
                     "These adjunctive therapies may benefit specific patient populations. Consider when standard therapy is insufficient.")
    
    # Smoking Cessation
    add_content_slide(prs, "Non-Pharmacological Management: Smoking Cessation",
                     """**Most Important Intervention!**
- Only intervention that slows disease progression
- Reduces mortality
- Improves symptoms

**Strategies:**
- Behavioral counseling
- Nicotine replacement therapy
- Varenicline (Chantix)
- Bupropion (Zyban)
- Combination approaches most effective

**5 A's Approach:**
1. Ask about tobacco use
2. Advise to quit
3. Assess willingness
4. Assist with quit plan
5. Arrange follow-up""",
                     "Smoking cessation is the only intervention proven to slow COPD progression. Every patient who smokes should be offered cessation support.")
    
    # Pulmonary Rehabilitation
    add_content_slide(prs, "Non-Pharmacological Management: Pulmonary Rehabilitation",
                     """**Components:**
- Exercise training
- Education
- Nutritional counseling
- Psychosocial support

**Benefits:**
- Improved exercise tolerance
- Reduced dyspnea
- Improved quality of life
- Reduced hospitalizations

**Program Structure:**
- 6-12 weeks, 2-3 sessions/week
- Maintenance program after completion
- Refer all symptomatic patients""",
                     "Pulmonary rehabilitation improves symptoms and quality of life. All symptomatic patients should be referred.")
    
    # Oxygen Therapy
    add_content_slide(prs, "Non-Pharmacological Management: Oxygen Therapy",
                     """**Indications (Long-term Oxygen Therapy - LTOT):**
- PaO2 ≤55 mmHg or SaO2 ≤88%
- PaO2 56-59 mmHg with cor pulmonale or polycythemia

**Goals:**
- Maintain SaO2 ≥90% during rest, sleep, and exertion
- Minimum 15 hours/day for mortality benefit

**Delivery Methods:**
- Nasal cannula (most common)
- Oxygen concentrators (home)
- Liquid oxygen (portable)
- Compressed gas (backup)

**Important:** Use caution with high-flow oxygen in CO2 retainers""",
                     "LTOT improves survival in hypoxemic patients. The 15-hour minimum is essential for mortality benefit.")
    
    # Vaccinations
    add_content_slide(prs, "Non-Pharmacological Management: Vaccinations",
                     """**Recommended Vaccines:**
- **Influenza:** Annual
- **Pneumococcal:**
  - PCV20 once, OR
  - PCV15 followed by PPSV23
- **COVID-19:** Per current guidelines
- **Tdap:** One-time, then Td booster

**Benefits:**
- Reduce exacerbations
- Reduce hospitalizations
- Reduce mortality""",
                     "Vaccinations are an essential part of COPD management to prevent exacerbations.")
    
    # Acute Exacerbations
    add_content_slide(prs, "Acute Exacerbations",
                     """**Definition:**
"An acute worsening of respiratory symptoms that results in additional therapy."

**Common Triggers:**
- Respiratory infections (viral or bacterial)
- Air pollution
- Non-adherence to therapy
- Heart failure
- Pulmonary embolism
- Pneumothorax

**Classification:**
- Mild: Treat with SABA only
- Moderate: Requires antibiotics and/or oral corticosteroids
- Severe: Requires hospitalization or ED visit""",
                     "Exacerbations are a major cause of morbidity and mortality. Prevention and prompt treatment are essential.")
    
    # Exacerbation Assessment
    add_content_slide(prs, "Exacerbation Management: Assessment",
                     """**History:**
- Baseline status vs. current
- Trigger identification
- Home medication use
- Prior intubations

**Physical Examination:**
- Work of breathing
- Ability to speak
- Mental status
- Cyanosis
- Accessory muscle use

**Diagnostic Testing:**
- ABG (if severe)
- Chest X-ray (rule out pneumonia, pneumothorax)
- ECG (rule out cardiac causes)
- Labs (CBC, BMP, sputum culture if indicated)""",
                     "Assessment determines severity and guides treatment. Look for complications like pneumonia or pneumothorax.")
    
    # Exacerbation Pharmacological Treatment
    add_content_slide(prs, "Exacerbation Management: Pharmacological Treatment",
                     """**Bronchodilators:**
- Short-acting beta-agonists (albuterol)
- Short-acting anticholinergics (ipratropium)
- Nebulizer or MDI with spacer (equivalent efficacy)

**Corticosteroids:**
- Prednisone 40mg daily for 5 days (evidence-based)
- No taper needed for short course

**Antibiotics (if bacterial infection suspected):**
- Amoxicillin-clavulanate
- Doxycycline
- Azithromycin
- Respiratory fluoroquinolone (severe)""",
 "Standard exacerbation treatment includes bronchodilators, corticosteroids, and antibiotics when indicated. The 5-day steroid course is evidence-based.")
    
    # Oxygen Therapy in Exacerbations
    add_content_slide(prs, "Exacerbation Management: Oxygen Therapy",
                     """**Goals:**
- Maintain SpO2 88-92%
- Avoid hyperoxia

**Rationale:**
- Chronic CO2 retainers have hypoxic drive
- High-flow oxygen can worsen hypercapnia
- Monitor for CO2 narcosis

**Monitoring:**
- Pulse oximetry
- ABG if concern for hypercapnia
- Mental status

**Warning Signs of CO2 Retention:**
- Decreased level of consciousness
- Headache
- Flapping tremor (asterixis)
- Acidosis on ABG""",
                     "Controlled oxygen therapy is essential in COPD exacerbations. Target 88-92% to avoid suppressing hypoxic drive.")
    
    # NIV
    add_content_slide(prs, "Exacerbation Management: Noninvasive Ventilation (NIV)",
                     """**Indications:**
- Acute respiratory acidosis (pH ≤7.35, PaCO2 ≥45)
- Severe dyspnea with clinical signs
- Persistent hypoxemia despite supplemental O2

**Benefits:**
- Reduces intubation rate
- Reduces mortality
- Reduces hospital LOS
- Allows respiratory muscle rest

**Contraindications:**
- Respiratory arrest
- Hemodynamic instability
- Inability to protect airway
- Facial trauma/surgery
- Uncooperative patient""",
                     "NIV is first-line therapy for COPD exacerbations with respiratory acidosis. It reduces intubation and mortality.")
    
    # NIV Implementation
    add_content_slide(prs, "NIV Implementation",
                     """**Initial Settings:**
- IPAP: 10 cmH2O (titrate to 15-20)
- EPAP: 5 cmH2O (titrate to 8-10)
- Backup rate: 10-12 if timed mode
- FIO2: Titrate to SpO2 88-92%

**Interface Selection:**
- Full-face mask (preferred initially)
- Nasal mask (if tolerated)
- Proper fit critical for success

**Monitoring:**
- Vital signs
- ABG at 1-2 hours
- Patient comfort
- Leaks
- Synchrony

**Success Criteria:**
- pH improvement within 2-4 hours
- Decreased respiratory rate
- Improved mental status""",
                     "Start with low pressures and titrate up. Monitor ABG at 1-2 hours to assess response.")
    
    # When to Escalate
    add_content_slide(prs, "When to Escalate Care",
                     """**Indications for ICU Admission:**
- Severe dyspnea unresponsive to initial therapy
- Altered mental status
- Persistent/worsening hypoxemia
- Severe acidosis (pH <7.25)
- Hemodynamic instability
- Respiratory arrest

**Indications for Intubation:**
- NIV failure (no improvement in 1-2 hours)
- Inability to protect airway
- Respiratory arrest
- Severe acidosis with deteriorating mental status
- Cardiac arrest""",
                     "Early recognition of NIV failure is critical. Don't delay intubation when indicated.")
    
    # Patient Education
    add_content_slide(prs, "Patient Education: Self-Management",
                     """**Action Plan Components:**
1. **Green Zone:** Doing well - continue routine medications
2. **Yellow Zone:** Symptoms worsening - increase bronchodilators, start oral steroids
3. **Red Zone:** Severe symptoms - seek emergency care

**Key Education Topics:**
- Proper inhaler technique
- Medication adherence
- Recognizing early signs of exacerbation
- When to seek help
- Smoking cessation
- Breathing techniques (pursed-lip, diaphragmatic)
- Energy conservation""",
                     "Patient education and self-management plans reduce exacerbations and improve outcomes.")
    
    # Inhaler Technique
    add_content_slide(prs, "Patient Education: Inhaler Technique",
                     """**Common Errors:**
- Not shaking MDI
- Poor coordination (MDI)
- Not exhaling fully before inhaling
- Not holding breath after inhalation
- Not rinsing mouth after ICS

**Teach-Back Method:**
1. Demonstrate technique
2. Patient demonstrates back
3. Correct errors
4. Provide written instructions
5. Reassess at each visit

**Spacer Use:**
- Improves drug delivery
- Reduces oropharyngeal deposition
- Essential for MDIs in acute setting""",
                     "Inhaler technique must be assessed at every visit. Spacers improve delivery and reduce side effects.")
    
    # Case Study 1
    add_content_slide(prs, "Case Study 1",
                     """**Patient:** Mr. Johnson, 68-year-old male
- 50 pack-year smoking history (quit 2 years ago)
- Diagnosis: COPD GOLD 3
- Current medications: Tiotropium, albuterol PRN

**Presentation:**
- Increased dyspnea for 3 days
- Increased sputum (yellow)
- Using albuterol every 2 hours

**Vital Signs:**
- HR: 102, BP: 148/88
- RR: 26, SpO2: 86% on room air
- Afebrile

**Questions:**
1. How would you classify this exacerbation?
2. What is your initial management?
3. What oxygen saturation target would you use?""",
                     "This is a typical moderate COPD exacerbation presentation.")
    
    # Case Study 1 Discussion
    add_content_slide(prs, "Case Study 1 Discussion",
                     """**Classification:**
- Moderate exacerbation (requires antibiotics and oral steroids)

**Initial Management:**
1. Controlled oxygen therapy (target SpO2 88-92%)
2. Nebulized bronchodilators (albuterol + ipratropium)
3. Oral corticosteroids (prednisone 40mg daily × 5 days)
4. Antibiotics (amoxicillin-clavulanate or doxycycline)
5. Consider NIV if acidosis present

**Oxygen Target:**
- 88-92% (COPD with likely CO2 retention)
- Monitor ABG for hypercapnia

**Disposition:**
- Likely requires hospitalization
- Home if good response to treatment and reliable follow-up""",
                     "Moderate exacerbations require antibiotics and steroids. Always use controlled oxygen in COPD patients.")
    
    # Case Study 2
    add_content_slide(prs, "Case Study 2",
                     """**Patient:** Mrs. Rodriguez, 72-year-old female
- COPD GOLD 4, on home oxygen 2L NC
- Admitted with severe exacerbation

**Current Status:**
- On 4L NC, SpO2 90%
- RR: 32, using accessory muscles
- Alert but anxious
- ABG: pH 7.28, PaCO2 68, PaO2 58, HCO3 32

**Questions:**
1. Interpret the ABG
2. Is this patient a candidate for NIV?
3. What NIV settings would you recommend?""",
                     "This patient has acute-on-chronic respiratory acidosis and is a candidate for NIV.")
    
    # Case Study 2 Discussion
    add_content_slide(prs, "Case Study 2 Discussion",
                     """**ABG Interpretation:**
- pH 7.28 = Acidemia
- PaCO2 68 = Respiratory acidosis
- HCO3 32 = Metabolic compensation (chronic)
- **Diagnosis:** Acute-on-chronic respiratory acidosis

**NIV Candidate:**
- YES - meets criteria:
  - Acute respiratory acidosis (pH ≤7.35)
  - Severe dyspnea
  - Persistent hypoxemia
  - Alert and cooperative

**NIV Settings:**
- IPAP: 10 cmH2O (titrate up)
- EPAP: 5 cmH2O
- FIO2: Titrate to SpO2 88-92%
- Full-face mask initially

**Monitoring:**
- Repeat ABG in 1-2 hours
- Watch for improvement in pH and RR""",
                     "This patient meets all criteria for NIV. Start with low pressures and titrate based on response.")
    
    # Key Takeaways
    add_content_slide(prs, "Key Takeaways",
                     """1. **COPD is Preventable and Treatable**
   - Smoking cessation is the most important intervention

2. **Diagnosis Requires Spirometry**
   - Post-bronchodilator FEV1/FVC <0.70

3. **Management is Multifaceted**
   - Pharmacological: Bronchodilators are mainstay
   - Non-pharmacological: Rehab, oxygen, vaccines

4. **Exacerbations Require Prompt Treatment**
   - Controlled oxygen (88-92%)
   - Bronchodilators, steroids, antibiotics
   - NIV for respiratory acidosis

5. **Patient Education is Essential**
   - Inhaler technique
   - Self-management action plans
   - Smoking cessation""",
                     "Review these key points. COPD management requires a comprehensive approach addressing all aspects of the disease.")
    
    # References
    add_content_slide(prs, "References",
                     """1. Global Initiative for Chronic Obstructive Lung Disease (GOLD). Global Strategy for the Diagnosis, Management, and Prevention of Chronic Obstructive Pulmonary Disease: 2024 Report.

2. Celli BR, et al. An Official American Thoracic Society/European Respiratory Society Statement: Research Questions in COPD. Am J Respir Crit Care Med. 2015;191(7):e4-e27.

3. Nici L, et al. American Thoracic Society/European Respiratory Society Statement on Pulmonary Rehabilitation. Am J Respir Crit Care Med. 2006;173(12):1390-1413.

4. Walters JA, et al. Different durations of corticosteroid therapy for exacerbations of chronic obstructive pulmonary disease. Cochrane Database Syst Rev. 2014;(12):CD006897.

5. AARC Clinical Practice Guideline: Oxygen Therapy for Adults in the Acute Care Facility. Respir Care. 2022.""",
                     "GOLD reports are updated annually and should be consulted for the most current recommendations.")
    
    add_title_slide(prs, "Thank You", "Respiratory Care Department\n\nImproving Lives Through Better Breathing")
    
    output_path = os.path.join(OUTPUT_DIR, "02-COPD-Management.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_copd_presentation()
