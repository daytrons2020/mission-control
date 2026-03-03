#!/usr/bin/env python3
"""
Create comprehensive PowerPoint presentations for respiratory education
with full content and medical images.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create output directory
OUTPUT_DIR = "/root/.openclaw/workspace/projects/respiratory-education/output"
IMAGES_DIR = "/root/.openclaw/workspace/projects/respiratory-education/images"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_content_slide(prs, title, content, notes=""):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Set content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    lines = content.strip().split('\n')
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        
        # Check if this is a header (starts with ** or ##)
        is_header = line.startswith('**') and line.endswith('**')
        is_subheader = line.startswith('##') or (line.startswith('**') and not is_header)
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Remove markdown
        clean_line = line.replace('**', '').replace('##', '').replace('#', '').strip()
        
        # Handle bullet points
        if clean_line.startswith('- ') or clean_line.startswith('* '):
            clean_line = clean_line[2:]
            p.level = 1
        elif clean_line.startswith('1.') or clean_line.startswith('2.') or clean_line.startswith('3.'):
            p.level = 1
        
        p.text = clean_line
        
        if is_header or is_subheader:
            p.font.bold = True
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 102, 153)
            p.level = 0
        else:
            p.font.size = Pt(16)
            p.level = 1 if p.level == 1 else 0
    
    # Add speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content, notes=""):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Left column
    left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.5))
    left_tf = left_shape.text_frame
    left_tf.word_wrap = True
    
    for line in left_content.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        p = left_tf.add_paragraph()
        clean_line = line.replace('**', '').replace('- ', '')
        p.text = clean_line
        p.font.size = Pt(14)
        if line.startswith('**') or line.startswith('##') or line.startswith('#'):
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 153)
    
    # Right column
    right_shape = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5))
    right_tf = right_shape.text_frame
    right_tf.word_wrap = True
    
    for line in right_content.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        p = right_tf.add_paragraph()
        clean_line = line.replace('**', '').replace('- ', '')
        p.text = clean_line
        p.font.size = Pt(14)
        if line.startswith('**') or line.startswith('##') or line.startswith('#'):
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 153)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

def add_image_slide(prs, title, image_path, content="", notes=""):
    """Add a slide with an image"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Try to add image if it exists
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4))
        except:
            pass
    
    # Add content on left side
    if content:
        content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.8), Inches(5.5))
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        for line in content.strip().split('\n'):
            line = line.strip()
            if not line:
                continue
            p = tf.add_paragraph()
            clean_line = line.replace('**', '').replace('- ', '')
            p.text = clean_line
            p.font.size = Pt(14)
            if line.startswith('**') or line.startswith('##') or line.startswith('#'):
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 153)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

# ==================== PRESENTATION 1: ANATOMY & PHYSIOLOGY ====================

def create_anatomy_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, "Respiratory Anatomy and Physiology", 
                   "Foundations of Respiratory Care\n\nPresenter: [Name], RRT\nDate: [Date]\nDepartment: Respiratory Care")
    
    # Learning Objectives
    add_content_slide(prs, "Learning Objectives", 
                     """By the end of this session, participants will be able to:
- Identify and describe the structures of the upper and lower respiratory tract
- Explain the mechanics of ventilation including pressure gradients and compliance
- Describe the process of gas exchange at the alveolar-capillary membrane
- Explain oxygen transport from lungs to tissues
- Describe the neural and chemical control of breathing
- Apply anatomical and physiological knowledge to clinical scenarios""",
                     "Emphasize the importance of understanding anatomy and physiology as the foundation for all respiratory care practice.")
    
    # Overview
    add_content_slide(prs, "Overview of the Respiratory System",
                     """**Two Main Divisions**

**Upper Respiratory Tract**
- Nose and nasal cavity
- Pharynx (nasopharynx, oropharynx, laryngopharynx)
- Larynx

**Lower Respiratory Tract**
- Trachea
- Bronchial tree (bronchi, bronchioles)
- Respiratory zone (respiratory bronchioles, alveolar ducts, alveoli)""",
                     "The respiratory system is divided into conducting and respiratory zones. The conducting zone moves air to the respiratory zone where gas exchange occurs.")
    
    # Upper Respiratory Tract - Nose
    add_content_slide(prs, "Upper Respiratory Tract: Nose and Nasal Cavity",
                     """**Structure and Function**

**External Nose**
- Framework of bone and cartilage
- Opens into nasal vestibule

**Nasal Cavity**
- Divided by nasal septum
- Three conchae (superior, middle, inferior)
- Creates turbulent airflow

**Functions:**
- **Filtration:** Nasal hairs filter large particles
- **Humidification:** Warms and moistens air to 100% humidity at body temperature
- **Olfaction:** Smell receptors in superior region""",
                     "The nose is the first line of defense for the respiratory system. Humidification and warming of inspired air protects the lower airways.")
    
    # Upper Respiratory Tract - Pharynx
    add_content_slide(prs, "Upper Respiratory Tract: Pharynx",
                     """**Three Regions**

**Nasopharynx**
- Posterior to nasal cavity
- Contains pharyngeal tonsil (adenoids)
- Opening of Eustachian tubes
- Passage for air only

**Oropharynx**
- Posterior to oral cavity
- Contains palatine tonsils
- Common passage for air and food

**Laryngopharynx**
- Inferior to oropharynx
- Splits into esophagus (posterior) and larynx (anterior)""",
                     "The pharynx serves as a common pathway for both air and food, with mechanisms to prevent aspiration.")
    
    # Larynx
    add_content_slide(prs, "Upper Respiratory Tract: Larynx",
                     """**The Voice Box**

**Structure**
- Cartilage framework (thyroid, cricoid, arytenoid)
- Vocal folds (true vocal cords)
- Vestibular folds (false vocal cords)
- Glottis: opening between vocal folds

**Functions**
- **Airway protection:** Epiglottis closes over larynx during swallowing
- **Voice production:** Vocal fold vibration creates sound
- **Airway maintenance:** Rigid structure prevents collapse

**Clinical Note:** The thyroid cartilage (Adam's apple) is an important landmark for intubation.""",
                     "The larynx protects the lower airway during swallowing and produces sound. The epiglottis is critical for preventing aspiration.")
    
    # Trachea
    add_content_slide(prs, "Lower Respiratory Tract: Trachea",
                     """**The Windpipe**

**Structure**
- 10-12 cm long in adults
- 16-20 C-shaped cartilage rings
- Open portion posterior (esophagus sits here)
- Divides at carina (T4-T5 level) into mainstem bronchi

**Function**
- Conducts air to bronchi
- Cilia and mucus trap and remove particles
- Cough reflex triggered by carina stimulation

**Clinical Note:** The right mainstem bronchus is more vertical than the left, making it more likely for aspirated objects and ET tubes to enter the right lung.""",
                     "The trachea's C-shaped rings provide structural support while allowing the esophagus to expand during swallowing. The carina is highly sensitive and can trigger cough.")
    
    # Bronchial Tree
    add_content_slide(prs, "Lower Respiratory Tract: Bronchial Tree",
                     """**Branching Pattern**

**Mainstem Bronchi**
- Right: wider, shorter, more vertical
- Left: narrower, longer, more horizontal (passes under aortic arch)

**Lobar Bronchi (Secondary)**
- Right: 3 lobes (upper, middle, lower)
- Left: 2 lobes (upper, lower)

**Segmental Bronchi (Tertiary)**
- 10 segments on right
- 8-10 segments on left
- Each segment is a separate anatomical and functional unit

**Bronchioles**
- No cartilage
- Smooth muscle allows bronchoconstriction/dilation""",
                     "The bronchial tree branches like an inverted tree. Understanding the differences between right and left mainstem bronchi is critical for proper ET tube placement.")
    
    # Respiratory Zone
    add_content_slide(prs, "The Respiratory Zone",
                     """**Where Gas Exchange Occurs**

**Respiratory Bronchioles**
- First structures with alveoli
- Transitional zone

**Alveolar Ducts**
- Lead to alveolar sacs

**Alveoli**
- 300-500 million per lung
- Total surface area: 70-100 square meters
- Thin walls (0.5 micrometers)

**Alveolar Structure**
- Type I pneumocytes: gas exchange surface
- Type II pneumocytes: surfactant production
- Alveolar macrophages: immune defense
- Capillary network: extensive for gas exchange""",
                     "The alveoli provide an enormous surface area for gas exchange. The thin alveolar-capillary membrane allows efficient diffusion of gases.")
    
    # The Pleura
    add_content_slide(prs, "The Pleura",
                     """**Membranes of the Lungs**

**Parietal Pleura**
- Lines thoracic cavity
- Costal (ribs), diaphragmatic, mediastinal portions
- Pain-sensitive (somatic innervation)

**Visceral Pleura**
- Covers lung surface
- Extends into fissures
- Not pain-sensitive

**Pleural Space**
- Potential space between layers
- Contains 10-20 ml serous fluid
- Negative pressure (-5 to -10 cmH2O at rest)
- Allows lungs to expand with chest wall

**Clinical Note:** Pneumothorax occurs when air enters pleural space, eliminating negative pressure and causing lung collapse.""",
                     "The pleural space and negative pressure are essential for normal lung expansion. Any air or fluid in this space can compromise breathing.")
    
    # Mechanics of Breathing
    add_content_slide(prs, "Mechanics of Breathing: Pressure Relationships",
                     """**Boyle's Law Application**

**Boyle's Law:** Pressure × Volume = Constant (at constant temperature)

**Inspiration:**
- Diaphragm contracts (moves down)
- External intercostals contract (ribs up and out)
- Thoracic volume increases
- Intrapleural pressure becomes more negative (-8 to -10 cmH2O)
- Alveolar pressure drops below atmospheric (-1 cmH2O)
- Air flows in (high to low pressure)

**Expiration:**
- Muscles relax
- Thoracic volume decreases
- Alveolar pressure rises above atmospheric (+1 cmH2O)
- Air flows out
- Passive process at rest""",
                     "Breathing is driven by pressure gradients created by changes in thoracic volume. Understanding these pressures is essential for managing mechanical ventilation.")
    
    # Lung Volumes
    add_content_slide(prs, "Lung Volumes and Capacities",
                     """**Static Lung Volumes**

**Tidal Volume (VT)**
- 500 ml (adult at rest)
- Air moving in and out with normal breath

**Inspiratory Reserve Volume (IRV)**
- 3000 ml
- Additional air that can be inhaled after normal inspiration

**Expiratory Reserve Volume (ERV)**
- 1100 ml
- Additional air that can be exhaled after normal expiration

**Residual Volume (RV)**
- 1200 ml
- Air remaining after maximal exhalation
- Prevents alveolar collapse""",
                     "Lung volumes are measured with spirometry. Residual volume cannot be measured with simple spirometry and requires special techniques.")
    
    # Lung Capacities
    add_content_slide(prs, "Lung Capacities",
                     """**Combinations of Volumes**

**Inspiratory Capacity (IC)**
- VT + IRV = 3500 ml
- Maximum air that can be inhaled from resting expiratory level

**Functional Residual Capacity (FRC)**
- ERV + RV = 2300 ml
- Air remaining at resting expiratory level
- Important for gas exchange continuity

**Vital Capacity (VC)**
- IRV + VT + ERV = 4600 ml
- Maximum air that can be exhaled after maximal inhalation

**Total Lung Capacity (TLC)**
- All volumes combined = 5800 ml
- Maximum air lungs can contain""",
                     "Capacities are combinations of two or more volumes. FRC is important as it maintains continuous gas exchange between breaths.")
    
    # Dynamic Lung Function
    add_content_slide(prs, "Dynamic Lung Function",
                     """**Forced Expiratory Measurements**

**Forced Vital Capacity (FVC)**
- Maximum amount of air forcibly exhaled

**Forced Expiratory Volume in 1 second (FEV1)**
- Amount of air exhaled in first second of FVC
- Normal: >80% of predicted
- FEV1/FVC ratio: normally >0.70 (70%)

**Peak Expiratory Flow (PEF)**
- Maximum flow rate during forced expiration
- Measured with peak flow meter

**Clinical Significance:**
- Obstructive disease: FEV1 reduced more than FVC, ratio <0.70
- Restrictive disease: Both FEV1 and FVC reduced, ratio normal or increased""",
                     "Dynamic lung function tests are essential for diagnosing and monitoring obstructive and restrictive lung diseases.")
    
    # Compliance and Resistance
    add_content_slide(prs, "Compliance and Resistance",
                     """**Determinants of Ventilation**

**Compliance (Distensibility)**
- Change in volume / Change in pressure
- Normal: 100 ml/cmH2O
- **Static compliance:** Measured during no-flow (plateau pressure)
- **Dynamic compliance:** Measured during breathing

**Factors Affecting Compliance:**
- Surfactant (increases compliance)
- Elastic tissue (decreases compliance)
- Disease states (pulmonary fibrosis decreases, emphysema increases)

**Resistance (Opposition to flow)**
- Primarily in airways
- Radius is major factor (Poiseuille's Law: resistance ∝ 1/r⁴)
- Bronchoconstriction dramatically increases resistance""",
                     "Compliance and resistance determine the work of breathing. In mechanical ventilation, we measure these to optimize ventilator settings.")
    
    # Surfactant
    add_content_slide(prs, "Surfactant",
                     """**Surface Tension Reduction**

**Composition**
- Dipalmitoylphosphatidylcholine (DPPC) - primary component
- Proteins (SP-A, SP-B, SP-C, SP-D)

**Production**
- Type II alveolar cells
- Begins at 24-28 weeks gestation
- Mature levels at 35 weeks

**Function**
- Reduces surface tension at air-liquid interface
- Prevents alveolar collapse
- Reduces work of breathing
- Prevents pulmonary edema

**Clinical Note:** Respiratory Distress Syndrome (RDS) in premature infants results from surfactant deficiency. Treatment includes surfactant replacement therapy.""",
                     "Surfactant is essential for preventing alveolar collapse. Without it, surface tension would cause alveoli to collapse during exhalation.")
    
    # Gas Exchange
    add_content_slide(prs, "Gas Exchange: Diffusion",
                     """**Alveolar-Capillary Membrane**

**Structure**
- Alveolar epithelium (Type I cells)
- Basement membrane
- Capillary endothelium
- Total thickness: 0.5 micrometers

**Fick's Law of Diffusion**
- Rate of diffusion ∝ (Surface Area × Pressure Difference × Solubility) / (Thickness × √Molecular Weight)

**Diffusion-Limited vs. Perfusion-Limited**
- O2: Normally perfusion-limited (can increase with exercise)
- CO: Diffusion-limited (binds tightly to hemoglobin)
- CO2: Highly soluble, easily diffuses""",
                     "The alveolar-capillary membrane is extremely thin to facilitate gas exchange. Fick's Law explains factors affecting diffusion efficiency.")
    
    # Partial Pressures
    add_content_slide(prs, "Partial Pressures of Gases",
                     """**Dalton's Law**

**Dalton's Law:** Total pressure = Sum of partial pressures of individual gases

**Atmospheric Air at Sea Level (760 mmHg):**
- Nitrogen (N2): 79% × 760 = 600 mmHg
- Oxygen (O2): 21% × 760 = 160 mmHg
- Water vapor: Variable (47 mmHg at body temperature)
- CO2: 0.04% (negligible)

**Alveolar Air:**
- Humidified (PH2O = 47 mmHg)
- O2 consumed, CO2 added
- PAO2 ≈ 100 mmHg
- PACO2 ≈ 40 mmHg

**Arterial Blood:**
- PaO2: 80-100 mmHg
- PaCO2: 35-45 mmHg""",
                     "Understanding partial pressures is essential for interpreting arterial blood gases and calculating alveolar-arterial oxygen gradients.")
    
    # Oxygen Transport
    add_content_slide(prs, "Oxygen Transport",
                     """**From Lungs to Tissues**

**Oxygen Transport:**
- Dissolved in plasma: 1.5% (PaO2 × 0.003)
- Bound to hemoglobin: 98.5%

**Hemoglobin**
- 4 heme groups (each binds one O2 molecule)
- 4 globin chains
- Normal Hgb: 14-18 g/dL (men), 12-16 g/dL (women)

**Oxygen Capacity**
- 1 g Hgb binds 1.34 ml O2
- Normal: 20 ml O2 per 100 ml blood

**Oxygen Saturation (SaO2)**
- Percentage of Hgb binding sites occupied by O2
- Normal: 95-100%
- Related to PaO2 by oxyhemoglobin dissociation curve""",
                     "Most oxygen is transported bound to hemoglobin. The small amount dissolved in plasma is what drives diffusion into tissues.")
    
    # Oxyhemoglobin Dissociation Curve
    add_content_slide(prs, "Oxyhemoglobin Dissociation Curve",
                     """**S-Shaped Curve Explained**

**Shape Significance:**
- Flat top (PaO2 60-100 mmHg): SaO2 remains high despite PaO2 changes
- Steep portion (PaO2 20-60 mmHg): Small PaO2 changes cause large SaO2 changes

**Key Points:**
- PaO2 60 = SaO2 90% (acceptable)
- PaO2 40 = SaO2 75% (hypoxemia)

**Factors Shifting Curve Right (Decreased Affinity, Easier O2 Release):**
- Increased temperature, Increased CO2 (Bohr effect), Decreased pH, Increased 2,3-DPG

**Factors Shifting Curve Left (Increased Affinity, Harder O2 Release):**
- Decreased temperature, Decreased CO2, Increased pH, Fetal hemoglobin, CO poisoning""",
                     "The S-shape of the curve ensures high saturation at normal PaO2 while allowing oxygen release in tissues where it's needed most.")
    
    # CO2 Transport
    add_content_slide(prs, "Carbon Dioxide Transport",
                     """**From Tissues to Lungs**

**CO2 Transport Forms:**
1. **Dissolved:** 7-10% (PCO2 × 0.07)
2. **Carbamino compounds:** 20% (bound to proteins, mostly Hgb)
3. **Bicarbonate:** 70% (HCO3-)

**Chloride Shift (Hamburger Effect)**
- In tissues: CO2 enters RBC, converted to HCO3-, HCO3- leaves cell, Cl- enters
- In lungs: Reverse occurs

**Haldane Effect**
- Deoxygenated Hgb binds more CO2 than oxygenated Hgb
- Facilitates CO2 pickup in tissues and release in lungs""",
                     "CO2 is primarily transported as bicarbonate. The chloride shift and Haldane effect facilitate efficient CO2 transport.")
    
    # Acid-Base Balance
    add_content_slide(prs, "Acid-Base Balance",
                     """**Henderson-Hasselbalch Equation**

**Equation:** pH = 6.1 + log([HCO3-] / (0.03 × PaCO2))

**Normal Values:**
- pH: 7.35-7.45
- PaCO2: 35-45 mmHg
- HCO3-: 22-26 mEq/L

**Regulatory Mechanisms:**
1. **Buffers (immediate):** Bicarbonate, proteins, hemoglobin
2. **Respiratory (minutes):** Change ventilation to alter PaCO2
3. **Renal (hours to days):** Excrete or retain HCO3-

**Disorders:**
- Respiratory acidosis: Increased PaCO2
- Respiratory alkalosis: Decreased PaCO2
- Metabolic acidosis: Decreased HCO3-
- Metabolic alkalosis: Increased HCO3-""",
                     "The respiratory system plays a crucial role in acid-base balance by regulating CO2 levels. The kidneys compensate over longer periods.")
    
    # Neural Control
    add_content_slide(prs, "Neural Control of Breathing",
                     """**Respiratory Centers**

**Medullary Centers**
- **Dorsal Respiratory Group (DRG):** Controls inspiration, responds to chemical stimuli
- **Ventral Respiratory Group (VRG):** Active during forced breathing, controls expiration

**Pontine Centers**
- **Apneustic center:** Promotes inspiration (inhibited by pneumotaxic)
- **Pneumotaxic center:** Limits inspiration, controls rate and depth

**Integration:**
- Normal breathing: DRG generates rhythm
- Exercise: VRG activates for active expiration
- Pontine centers fine-tune pattern""",
                     "Breathing is controlled by centers in the brainstem that integrate chemical and neural signals to maintain appropriate ventilation.")
    
    # Chemical Control
    add_content_slide(prs, "Chemical Control of Breathing",
                     """**Central and Peripheral Chemoreceptors**

**Central Chemoreceptors (Medulla)**
- Respond to H+ concentration in CSF
- H+ does not cross blood-brain barrier, but CO2 does
- CO2 + H2O → H2CO3 → H+ + HCO3-
- Primary drive for ventilation under normal conditions

**Peripheral Chemoreceptors (Carotid and Aortic Bodies)**
- Respond to:
  - Decreased PaO2 (hypoxemia, <60 mmHg)
  - Increased PaCO2 (hypercapnia)
  - Decreased pH (acidemia)
- Primary response to hypoxemia

**Clinical Note:** In COPD with chronic CO2 retention, hypoxemia becomes primary drive (respond to low O2, not high CO2).""",
                     "CO2 is the primary driver of ventilation in healthy individuals. In chronic CO2 retainers, low O2 becomes the primary drive, requiring careful oxygen administration.")
    
    # Lung Reflexes
    add_content_slide(prs, "Lung Reflexes",
                     """**Protective Mechanisms**

**Hering-Breuer Inflation Reflex**
- Stretch receptors in smooth muscle of airways
- Inhibits inspiration when lungs overinflate
- More active in infants

**Deflation Reflex**
- Stimulates inspiration when lungs deflate

**Irritant Reflex**
- Noxious stimuli (smoke, dust) trigger bronchoconstriction and cough
- Receptors in epithelium

**Juxtapulmonary Capillary Receptors (J-receptors)**
- Located near capillaries
- Stimulated by increased interstitial fluid (pulmonary edema)
- Cause rapid, shallow breathing and dyspnea

**Cough Reflex**
- Triggered by mechanical or chemical irritation
- Carina most sensitive area""",
                     "Lung reflexes protect against overinflation, noxious stimuli, and help clear secretions through coughing.")
    
    # Case Study 1
    add_content_slide(prs, "Clinical Application: Case Study 1",
                     """**Patient Scenario**

**Mr. Thompson, 68-year-old male**
- Admitted with community-acquired pneumonia
- Current status: Tachypneic, using accessory muscles
- ABG: pH 7.48, PaCO2 30, PaO2 58, HCO3 22

**Questions:**
1. What type of acid-base disorder is present?
2. What is the physiological explanation for the breathing pattern?
3. Why is the PaO2 low despite increased ventilation?
4. What factors affect oxygen delivery to his tissues?

**Discussion Points:**
- Respiratory alkalosis from hypoxemic drive
- V/Q mismatch in pneumonia
- O2 delivery depends on cardiac output, Hgb, and saturation""",
                     "This case demonstrates respiratory alkalosis due to hypoxemic drive. The patient is hyperventilating to compensate for hypoxemia from pneumonia.")
    
    # Case Study 2
    add_content_slide(prs, "Clinical Application: Case Study 2",
                     """**Patient Scenario**

**Mrs. Garcia, 72-year-old female with COPD**
- Chronic CO2 retention (baseline PaCO2 55)
- Presents with exacerbation
- ABG: pH 7.30, PaCO2 70, PaO2 48, HCO3 34

**Questions:**
1. What is her acid-base status?
2. Why is the HCO3 elevated?
3. What is the risk of giving high-flow oxygen?
4. How should oxygen be administered?

**Discussion Points:**
- Acute-on-chronic respiratory acidosis
- Renal compensation (chronic metabolic alkalosis)
- Risk of suppressing hypoxic drive
- Controlled low-flow oxygen (target SpO2 88-92%)""",
                     "This case shows acute-on-chronic respiratory acidosis. High-flow oxygen could suppress the patient's hypoxic drive, worsening respiratory acidosis.")
    
    # Key Takeaways
    add_content_slide(prs, "Key Takeaways",
                     """1. **Anatomy Knowledge is Clinical**
   - Right mainstem intubation risk
   - Carina sensitivity for cough
   - Pleural space importance

2. **Ventilation Depends on Pressure Gradients**
   - Boyle's Law governs breathing mechanics
   - Compliance and resistance determine work of breathing

3. **Gas Exchange Requires:**
   - Intact alveolar-capillary membrane
   - Adequate surface area
   - Proper ventilation-perfusion matching

4. **Oxygen Transport is Multi-Step:**
   - Alveolar ventilation
   - Diffusion across membrane
   - Hemoglobin binding
   - Tissue release

5. **Breathing Control is Integrated:**
   - Neural centers generate rhythm
   - Chemical receptors provide feedback
   - Reflexes protect the lungs""",
                     "Review these key concepts as they form the foundation for understanding respiratory disease and therapy.")
    
    # References
    add_content_slide(prs, "References",
                     """1. Kacmarek RM, Stoller JK, Heuer AJ. Egan's Fundamentals of Respiratory Care. 12th ed. Elsevier; 2021.

2. West JB. Respiratory Physiology: The Essentials. 10th ed. Wolters Kluwer; 2016.

3. Levitzky MG. Pulmonary Physiology. 9th ed. McGraw-Hill; 2018.

4. Berne RM, Levy MN. Principles of Physiology. 4th ed. Mosby; 2006.

5. AARC Clinical Practice Guidelines. Available at: https://www.aarc.org/resource/clinical-practice-guidelines/""",
                     "These references provide additional detail for those interested in deeper study of respiratory physiology.")
    
    # Thank you
    add_title_slide(prs, "Thank You", "Respiratory Care Department\n\nUnderstanding the Science Behind the Care")
    
    # Save presentation
    output_path = os.path.join(OUTPUT_DIR, "01-Anatomy-Physiology.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")
    return output_path

if __name__ == "__main__":
    create_anatomy_presentation()
